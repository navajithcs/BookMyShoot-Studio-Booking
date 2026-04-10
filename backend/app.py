from flask import Flask, render_template, redirect, url_for, flash, request, jsonify, send_file
from flask_login import login_user, logout_user, login_required, current_user
from extensions import db, login_manager, csrf, cors
from config import DATABASE_URI
from models import User, Photographer, Booking, Package, PortfolioItem, Notification, ActivityLog, Payment, EventCategory, ContactMessage, Feedback, ServicePackage, RefundRequest, RefundTransaction, SupportTicket
from admin import admin_bp
from photographer import photographer_bp
from customer import customer_bp
from invoice_generator import generate_booking_invoice_pdf
from email_service import send_booking_accepted_email
from validators import (
    validate_registration, validate_login, validate_booking, validate_contact,
    validate_profile_update, validate_photographer_update, validate_change_password,
    validate_feedback, validate_admin_create_user, validate_booking_status,
    sanitize_data, check_booking_conflicts, get_booked_slots,
    validate_numeric, validate_text_length, validate_enum, VALID_SERVICE_TYPES
)
from functools import wraps
import hmac
import hashlib
import os
import sys
import uuid
import signal
import socket
import logging
from datetime import datetime
from sqlalchemy import text
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

# ─── Logging ───
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'super_secret_booking_key_for_flash')

# Database Configuration
app.config['SQLALCHEMY_DATABASE_URI'] = DATABASE_URI
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['WTF_CSRF_CHECK_DEFAULT'] = False

# Initialize Extensions
db.init_app(app)
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message_category = 'info'
csrf.init_app(app)
<<<<<<< HEAD
import re
cors.init_app(app, resources={r"/api/*": {"origins": re.compile(r"https?://(localhost|127\.0\.0\.1)(:\d+)?")}}, supports_credentials=True)
=======
cors.init_app(app, resources={r"/api/*": {"origins": ["http://localhost:5173", "http://127.0.0.1:5173", "http://localhost:5000", "http://127.0.0.1:5000"]}}, supports_credentials=True)


>>>>>>> fba2361 (Initial commit)
@app.before_request
def enforce_csrf_for_non_api_routes():
    if request.path.startswith('/api/'):
        return
    if request.method in ('POST', 'PUT', 'PATCH', 'DELETE'):
        csrf.protect()

# Razorpay Configuration (optional)
RAZORPAY_KEY_ID = os.getenv("RAZORPAY_KEY_ID", "")
RAZORPAY_KEY_SECRET = os.getenv("RAZORPAY_KEY_SECRET", "")
razorpay_client = None
try:
    import razorpay
    if RAZORPAY_KEY_ID and RAZORPAY_KEY_SECRET:
        razorpay_client = razorpay.Client(auth=(RAZORPAY_KEY_ID, RAZORPAY_KEY_SECRET))
except ImportError:
    logger.warning("razorpay package not installed – payment features disabled")

# Register Blueprints
app.register_blueprint(admin_bp, url_prefix='/admin')
app.register_blueprint(photographer_bp, url_prefix='/photographer')
app.register_blueprint(customer_bp, url_prefix='/customer')

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or current_user.user_type != 'admin':
            flash('Admin access required.', 'error')
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def log_activity(user_id, action, details=None):
    """Helper to log admin/system activities"""
    try:
        log = ActivityLog(user_id=user_id, action=action, details=details)
        db.session.add(log)
        db.session.commit()
    except Exception as e:
        print(f"Error logging activity: {e}")

_BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PPT_FOLDER = os.path.join(_BASE_DIR, "generated_ppts")
os.makedirs(PPT_FOLDER, exist_ok=True)
UPLOAD_FOLDER = os.path.join(_BASE_DIR, "uploads", "portfolio")
PROFILE_FOLDER = os.path.join(_BASE_DIR, "uploads", "profiles")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROFILE_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['PROFILE_FOLDER'] = PROFILE_FOLDER

# Create database tables
with app.app_context():
    db.create_all()

    # Backward-compatible sqlite schema patching for local/dev DBs.
    # Prevents 500s when older DBs miss newly added columns.
    try:
        if db.engine.url.drivername.startswith('sqlite'):
            required_columns = {
                'users': {
                    'phone': "VARCHAR(20)",
                    'profile_image': "VARCHAR(255)",
                    'is_active': "BOOLEAN DEFAULT 1",
                    'updated_at': "DATETIME"
                },
                'photographers': {
                    'experience': "VARCHAR(50)",
                    'location': "VARCHAR(255)"
                },
                'bookings': {
                    'package_name': "VARCHAR(100)"
                }
            }

            for table_name, columns in required_columns.items():
                existing_columns = {
                    row[1]
                    for row in db.session.execute(text(f"PRAGMA table_info({table_name})")).fetchall()
                }
                for column_name, column_type in columns.items():
                    if column_name not in existing_columns:
                        db.session.execute(
                            text(f"ALTER TABLE {table_name} ADD COLUMN {column_name} {column_type}")
                        )

            db.session.commit()
    except Exception as e:
        db.session.rollback()
        print(f"Schema patch warning: {e}")


@app.route('/')
def home():
    return render_template('index.html')

@app.route('/contact')
def contact():
    return render_template('contact.html')

@app.route('/services')
def services():
    # A generic services overview page or redirect to home section
    return redirect(url_for('home', _anchor='services'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        if current_user.user_type == 'admin':
            return redirect(url_for('admin.dashboard'))
        elif current_user.user_type == 'photographer':
            return redirect(url_for('photographer.dashboard'))
        return redirect(url_for('customer.dashboard'))

    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        user = User.query.filter_by(email=email).first()
        
        if user and user.check_password(password):
            login_user(user)
            log_activity(user.id, 'login', f'User logged in via SSR: {email}')
            flash(f'Welcome back, {user.first_name}!', 'success')
            
            if user.user_type == 'admin':
                return redirect(url_for('admin.dashboard'))
            elif user.user_type == 'photographer':
                return redirect(url_for('photographer.dashboard'))
            return redirect(url_for('customer.dashboard'))
        
        flash('Invalid email or password.', 'error')
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        first_name = request.form.get('first_name')
        last_name = request.form.get('last_name')
        email = request.form.get('email')
        password = request.form.get('password')
        user_type = request.form.get('user_type', 'customer')
        
        if User.query.filter_by(email=email).first():
            flash('Email already registered.', 'error')
            return redirect(url_for('register'))
            
        user = User(first_name=first_name, last_name=last_name, email=email, user_type=user_type)
        user.set_password(password)
        db.session.add(user)
        db.session.commit()
        
        if user_type == 'photographer':
            photographer = Photographer(user_id=user.id)
            db.session.add(photographer)
            db.session.commit()
            
        flash('Account created successfully! Please login.', 'success')
        return redirect(url_for('login'))
        
    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    flash('Logged out successfully.', 'success')
    return redirect(url_for('login'))

@app.route('/uploads/profiles/<filename>')
def uploaded_profile_file(filename):
    return send_file(os.path.join(app.config['PROFILE_FOLDER'], filename))

# Booking logic moved to Blueprints.

# ==================== AUTH ROUTES ====================

@app.route('/api/register', methods=['POST'])
def api_register():
    """Register a new user"""
    try:
        data = sanitize_data(request.json or {})
        
        # Validate all fields
        errors = validate_registration(data)
        if errors:
            return jsonify({'error': errors[0], 'errors': errors}), 400
        
        # Check if user already exists
        if User.query.filter_by(email=data['email']).first():
            return jsonify({'error': 'Email already registered'}), 400
        
        # Create new user
        user = User(
            first_name=data['first_name'],
            last_name=data['last_name'],
            email=data['email'],
            phone=data.get('phone', ''),
            user_type=data['user_type']
        )
        user.set_password(data['password'])
        
        db.session.add(user)
        db.session.flush() # Get user.id without committing
        
        # If photographer, create photographer profile
        if data['user_type'] == 'photographer':
            photographer = Photographer(
                user_id=user.id,
                specialty=data.get('specialty', 'General'), # Default specialty
                hourly_rate=data.get('hourly_rate', 0),
                bio=data.get('bio', '')
            )
            db.session.add(photographer)
            
        db.session.commit()
        
        # If photographer, add photographer_id to response
        if data['user_type'] == 'photographer':
            user_data = user.to_dict()
            user_data['photographer_id'] = photographer.id
            
            return jsonify({
                'message': 'Registration successful',
                'user': user_data
            }), 201
        
        return jsonify({
            'message': 'Registration successful',
            'user': user.to_dict()
        }), 201
    except Exception as e:
        # Return full error message in development to help debugging
        return jsonify({'error': f'Registration error: {e}'}), 500


@app.route('/api/login', methods=['POST'])
def api_login():
    """Login user"""
    data = sanitize_data(request.json or {})
    
    # Validate input
    errors = validate_login(data)
    if errors:
        return jsonify({'error': errors[0]}), 400
    
    email = data.get('email')
    password = data.get('password')
    
    # Find user by email
    user = User.query.filter_by(email=email).first()
    
    if not user or not user.check_password(password):
        return jsonify({'error': 'Invalid email or password'}), 401
    
    if not user.is_active:
        return jsonify({'error': 'Account is deactivated'}), 403
    
    # Log successful login
    log_activity(user.id, 'login', f'User logged in: {email}')
    
    user_data = user.to_dict()
    
    # If photographer, append photographer_id
    if user.user_type == 'photographer':
        photographer = Photographer.query.filter_by(user_id=user.id).first()
        if photographer:
            user_data['photographer_id'] = photographer.id
    
    return jsonify({
        'message': 'Login successful',
        'user': user_data
    }), 200


@app.route('/api/logout', methods=['POST'])
def api_logout():
    """Logout user"""
    data = request.json
    user_id = data.get('user_id')
    
    if user_id:
        user = User.query.get(user_id)
        if user:
            log_activity(user_id, 'logout', f'User logged out: {user.email}')
    
    return jsonify({'message': 'Logged out successfully'}), 200


@app.route('/api/user/update-profile', methods=['POST'])
def update_profile():
    """Update user profile"""
    data = sanitize_data(request.json or {})
    user_id = data.get('user_id')
    
    if not user_id:
        return jsonify({'error': 'User ID required'}), 400
    
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    # Validate fields
    errors = validate_profile_update(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400
    
    first_name = data.get('first_name')
    last_name = data.get('last_name')
    phone = data.get('phone')
    email = data.get('email')
    
    # Check if email is being changed and if new email already exists
    if email and email != user.email:
        existing = User.query.filter_by(email=email).first()
        if existing:
            return jsonify({'error': 'Email already in use'}), 400
        user.email = email
    
    if first_name:
        user.first_name = first_name
    if last_name:
        user.last_name = last_name
    if phone is not None:
        user.phone = phone
    
    db.session.commit()
    
    return jsonify({
        'message': 'Profile updated successfully',
        'user': user.to_dict()
    }), 200


@app.route('/api/user/change-password', methods=['POST'])
def change_password():
    """Change user password"""
    data = request.json or {}
    user_id = data.get('user_id')
    current_password = data.get('current_password')
    new_password = data.get('new_password')
    
    if not user_id or not current_password or not new_password:
        return jsonify({'error': 'All fields are required'}), 400
    
    # Validate new password strength
    errors = validate_change_password(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400
    
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if not user.check_password(current_password):
        return jsonify({'error': 'Current password is incorrect'}), 400
    
    user.set_password(new_password)
    db.session.commit()
    
    return jsonify({'message': 'Password changed successfully'}), 200


@app.route('/api/user/<int:user_id>', methods=['GET'])
def get_user(user_id):
    """Get user by ID"""
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    return jsonify({'user': user.to_dict()}), 200


@app.route('/api/user/upload-avatar', methods=['POST'])
def upload_avatar():
    """Upload user profile image"""
    user_id = request.form.get('user_id')
    if not user_id:
        return jsonify({'error': 'User ID required'}), 400
    
    user = User.query.get(int(user_id))
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if 'avatar' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['avatar']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    allowed = {'jpg', 'jpeg', 'png', 'gif', 'webp'}
    ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
    if ext not in allowed:
        return jsonify({'error': 'Invalid file type. Allowed: jpg, jpeg, png, gif, webp'}), 400
    
    import uuid
    filename = f"user_{user_id}_{uuid.uuid4().hex[:8]}.{ext}"
    filepath = os.path.join(PROFILE_FOLDER, filename)
    
    # Remove old profile image if exists
    if user.profile_image:
        old_path = os.path.join(PROFILE_FOLDER, user.profile_image)
        if os.path.exists(old_path):
            os.remove(old_path)
    
    file.save(filepath)
    user.profile_image = filename
    db.session.commit()
    
    return jsonify({
        'message': 'Profile image updated',
        'user': user.to_dict(),
        'image_url': f'/uploads/profiles/{filename}'
    }), 200


@app.route('/api/user/remove-avatar', methods=['POST'])
def remove_avatar():
    """Remove user profile image"""
    data = request.json
    user_id = data.get('user_id')
    if not user_id:
        return jsonify({'error': 'User ID required'}), 400

    user = User.query.get(int(user_id))
    if not user:
        return jsonify({'error': 'User not found'}), 404

    if user.profile_image:
        old_path = os.path.join(PROFILE_FOLDER, user.profile_image)
        if os.path.exists(old_path):
            os.remove(old_path)
        user.profile_image = None
        db.session.commit()

    return jsonify({
        'message': 'Profile image removed',
        'user': user.to_dict()
    }), 200


@app.route('/api/user/delete-account', methods=['POST'])
def delete_own_account():
    """Allow a user to delete their own account"""
    data = request.json
    user_id = data.get('user_id')
    password = data.get('password')
    
    if not user_id or not password:
        return jsonify({'error': 'User ID and password are required'}), 400
    
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if not user.check_password(password):
        return jsonify({'error': 'Incorrect password'}), 400
    
    try:
        from models import Booking, Notification, ActivityLog
        Booking.query.filter_by(customer_id=user_id).delete()
        Notification.query.filter_by(user_id=user_id).delete()
        ActivityLog.query.filter_by(user_id=user_id).delete()
        
        if user.profile_image:
            old_path = os.path.join(PROFILE_FOLDER, user.profile_image)
            if os.path.exists(old_path):
                os.remove(old_path)
        
        db.session.delete(user)
        db.session.commit()
        return jsonify({'message': 'Account deleted successfully'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


# ==================== BOOKING ROUTES ====================

@app.route('/api/bookings', methods=['POST'])
def create_booking():
    """Create a new booking"""
    data = sanitize_data(request.json or {})
    
    # Validate all booking fields
    errors = validate_booking(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400
    
    # Verify customer exists
    customer = User.query.get(data['customer_id'])
    if not customer or customer.user_type != 'customer':
        return jsonify({'error': 'Invalid customer'}), 400
    
    # If photographer_id is provided, verify photographer exists
    if data.get('photographer_id'):
        photographer = Photographer.query.get(data['photographer_id'])
        if not photographer:
            return jsonify({'error': 'Invalid photographer'}), 400
        
        # Check for booking conflicts
        conflict = check_booking_conflicts(
            data['photographer_id'],
            data['event_date'],
            data.get('event_time')
        )
        if conflict:
            return jsonify({'error': conflict}), 409
    
    # Calculate token and remaining amounts (20% token)
    total_price = float(data['total_price'])
    token_amount = total_price * 0.20
    remaining_amount = total_price - token_amount
    
    # Parse date
    try:
        event_date = datetime.strptime(data['event_date'], '%Y-%m-%d').date()
    except ValueError:
        return jsonify({'error': 'Invalid date format. Use YYYY-MM-DD'}), 400
    
    # Build full location string from structured address fields
    address_parts = [p for p in [data.get('address_line', ''), data.get('state', ''), data.get('pincode', ''), data.get('country', '')] if p]
    full_location = ', '.join(address_parts) if address_parts else data.get('location', '')

    # Create booking
    booking = Booking(
        customer_id=data['customer_id'],
        photographer_id=data.get('photographer_id'),  # Can be None for open requests
        service_type=data['service_type'],
        package_name=data.get('package_name', ''),
        event_date=event_date,
        event_time=data.get('event_time', ''),
        location=full_location,
        location_type=data.get('location_type', ''),
        address_line=data.get('address_line', ''),
        pincode=data.get('pincode', ''),
        state=data.get('state', ''),
        country=data.get('country', 'India'),
        contact_number=data.get('contact_number', ''),
        notes=data.get('notes', ''),
        total_price=total_price,
        token_amount=token_amount,
        remaining_amount=remaining_amount,
        status='pending',
        payment_status='pending'
    )
    
    db.session.add(booking)
    db.session.commit()
    
    # Send confirmation notification to the customer
    customer_name = f"{customer.first_name}" if customer else 'Customer'
    customer_notification = Notification(
        user_id=data['customer_id'],
        booking_id=booking.id,
        title='Booking Submitted! 📸',
        message=f'Your {booking.service_type} booking for {booking.event_date} has been submitted. {"The photographer will review and respond soon." if data.get("photographer_id") else "Available photographers can accept your request."}',
        notification_type='booking_created',
        is_read=False
    )
    db.session.add(customer_notification)
    
    # If a specific photographer was selected, send them a notification
    if data.get('photographer_id'):
        photographer = Photographer.query.get(data['photographer_id'])
        if photographer:
            customer = User.query.get(data['customer_id'])
            customer_name = f"{customer.first_name} {customer.last_name}" if customer else 'A customer'
            
            notification = Notification(
                user_id=photographer.user_id,
                booking_id=booking.id,
                title='New Booking Request! 📸',
                message=f'{customer_name} has requested you for {booking.service_type} on {booking.event_date}. View details to accept or decline.',
                notification_type='new_booking_request',
                is_read=False
            )
            db.session.add(notification)
    
    db.session.commit()
    
    return jsonify({
        'message': 'Booking created successfully',
        'booking': booking.to_dict()
    }), 201


@app.route('/api/bookings/booked-slots', methods=['GET'])
def get_booked_slots_api():
    """Get booked time slots for a photographer on a given date"""
    photographer_id = request.args.get('photographer_id')
    event_date = request.args.get('event_date')
    if not photographer_id or not event_date:
        return jsonify({'error': 'photographer_id and event_date are required'}), 400
    try:
        photographer_id = int(photographer_id)
    except (TypeError, ValueError):
        return jsonify({'error': 'Invalid photographer_id'}), 400
    slots = get_booked_slots(photographer_id, event_date)
    return jsonify({'booked_slots': slots, 'date': event_date, 'photographer_id': photographer_id}), 200


@app.route('/api/bookings/requests', methods=['GET'])
def get_booking_requests():
    """Get all pending booking requests (without assigned photographer)"""
    bookings = Booking.query.filter_by(photographer_id=None, status='pending').order_by(Booking.created_at.desc()).all()
    return jsonify({
        'bookings': [b.to_dict() for b in bookings]
    }), 200


@app.route('/api/bookings/requests/photographer/<int:photographer_id>', methods=['GET'])
def get_filtered_booking_requests(photographer_id):
    """Get pending booking requests filtered by photographer's specialty"""
    photographer = Photographer.query.get(photographer_id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404
    
    # Get photographer's specialty (can be comma-separated like 'Wedding,Engagement,Events')
    specialty = photographer.specialty
    
    # Map specialty keywords to the service_type values used in bookings
    specialty_to_services = {
        'wedding': ['wedding'],
        'engagement': ['engagement'],
        'birthday': ['birthday'],
        'baby shower': ['babyshower'],
        'babyshower': ['babyshower'],
        'naming': ['naming'],
        'studio': ['studio'],
        'portrait': ['studio', 'portrait'],
        'events': ['birthday', 'babyshower', 'naming', 'engagement', 'wedding'],
        'event': ['birthday', 'babyshower', 'naming', 'engagement', 'wedding'],
        'photography': ['wedding', 'engagement', 'birthday', 'babyshower', 'naming', 'studio'],
    }
    
    # Build list of service types from photographer's specialties
    matching_services = []
    if specialty:
        # Split comma-separated specialties and check each one
        parts = [s.strip().lower() for s in specialty.split(',')]
        for part in parts:
            # Check direct mapping
            if part in specialty_to_services:
                matching_services.extend(specialty_to_services[part])
            else:
                # Check if any keyword is contained in the specialty part
                for keyword, services in specialty_to_services.items():
                    if keyword in part:
                        matching_services.extend(services)
            # Also add the part itself as a service type
            matching_services.append(part)
        matching_services = list(set(matching_services))  # Remove duplicates
    
    # Query for open bookings (no photographer assigned yet)
    if matching_services:
        bookings = Booking.query.filter(
            Booking.photographer_id == None,
            Booking.status == 'pending',
            Booking.service_type.in_(matching_services)
        ).order_by(Booking.created_at.desc()).all()
    else:
        # No specialty set — show all open requests
        bookings = Booking.query.filter_by(photographer_id=None, status='pending').order_by(Booking.created_at.desc()).all()
    
    return jsonify({
        'bookings': [b.to_dict() for b in bookings]
    }), 200


@app.route('/api/bookings/<int:booking_id>/accept', methods=['PUT'])
def accept_booking_request(booking_id):
    """Accept a booking request (photographer claims the job)"""
    data = request.json
    photographer_id = data.get('photographer_id')
    
    if not photographer_id:
        return jsonify({'error': 'Photographer ID is required'}), 400
    
    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404
    
    # For direct requests: photographer_id is already set to THIS photographer
    # For open requests: photographer_id is None
    if booking.photographer_id is not None and booking.photographer_id != photographer_id:
        return jsonify({'error': 'This booking has already been accepted by another photographer'}), 400
    
    if booking.status != 'pending':
        return jsonify({'error': 'This booking is no longer available'}), 400
    
    # Verify photographer exists
    photographer = Photographer.query.get(photographer_id)
    if not photographer:
        return jsonify({'error': 'Invalid photographer'}), 400
    
    # Check for time-slot conflicts before accepting
    conflict = check_booking_conflicts(
        photographer_id,
        str(booking.event_date),
        booking.event_time,
        exclude_booking_id=booking_id
    )
    if conflict:
        return jsonify({'error': conflict}), 409
    
    # Assign photographer and update status
    booking.photographer_id = photographer_id
    booking.status = 'accepted'
    
    # Notify the customer that their booking was accepted
    customer = User.query.get(booking.customer_id)
    photographer_name = f"{photographer.user.first_name} {photographer.user.last_name}" if photographer.user else 'A Photographer'
    
    notification = Notification(
        user_id=booking.customer_id,
        booking_id=booking.id,
        title='Booking Accepted! \U0001f389',
        message=f'{photographer_name} has accepted your {booking.service_type} booking for {booking.event_date}. Please pay the token amount from My Orders to confirm.',
        notification_type='job_accepted',
        is_read=False
    )
    db.session.add(notification)
    db.session.commit()
    
    # Send email with PDF invoice to customer
    try:
        if customer and photographer:
            pdf_bytes = generate_booking_invoice_pdf(booking, customer, photographer)
            send_booking_accepted_email(booking, customer, photographer, pdf_bytes)
    except Exception as e:
        print(f"[EMAIL ERROR] Failed to send booking acceptance email: {e}")
    
    return jsonify({
        'message': 'Booking accepted successfully',
        'booking': booking.to_dict()
    }), 200


@app.route('/api/bookings/<int:booking_id>', methods=['GET'])
def get_booking(booking_id):
    """Get booking by ID"""
    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404
    
    return jsonify({'booking': booking.to_dict()}), 200


@app.route('/api/bookings/<int:booking_id>/invoice', methods=['GET'])
def download_booking_invoice(booking_id):
    """Download PDF invoice for a booking"""
    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404
    
    customer = User.query.get(booking.customer_id)
    photographer = Photographer.query.get(booking.photographer_id)
    
    if not customer or not photographer:
        return jsonify({'error': 'Booking data incomplete'}), 400
    
    try:
        from io import BytesIO
        pdf_bytes = generate_booking_invoice_pdf(booking, customer, photographer)
        pdf_buffer = BytesIO(pdf_bytes)
        pdf_buffer.seek(0)
        return send_file(
            pdf_buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f'BookMyShoot_Invoice_{booking.id}.pdf'
        )
    except Exception as e:
        return jsonify({'error': f'Failed to generate invoice: {str(e)}'}), 500


@app.route('/api/bookings/customer/<int:customer_id>', methods=['GET'])
def get_customer_bookings(customer_id):
    """Get all bookings for a customer"""
    bookings = Booking.query.filter_by(customer_id=customer_id).order_by(Booking.created_at.desc()).all()
    return jsonify({
        'bookings': [b.to_dict() for b in bookings]
    }), 200


@app.route('/api/bookings/photographer/<int:photographer_id>', methods=['GET'])
def get_photographer_bookings(photographer_id):
    """Get all bookings for a photographer"""
    bookings = Booking.query.filter_by(photographer_id=photographer_id).order_by(Booking.created_at.desc()).all()
    return jsonify({
        'bookings': [b.to_dict() for b in bookings]
    }), 200


@app.route('/api/bookings/<int:booking_id>/status', methods=['PUT'])
def update_booking_status(booking_id):
    """Update booking status (accept/decline/complete)"""
    data = request.json
    booking = Booking.query.get(booking_id)
    
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404
    
    old_status = booking.status
    
    if 'status' in data:
        # Validate status value
        err = validate_booking_status(data['status'])
        if err:
            return jsonify({'error': err}), 400
        booking.status = data['status']
    
    if 'photographer_id' in data:
        booking.photographer_id = data['photographer_id']
    
    # Create notification for customer when photographer accepts/declines
    if data.get('status') == 'accepted' and old_status != 'accepted':
        photographer = Photographer.query.get(booking.photographer_id)
        photographer_name = photographer.user.first_name if photographer and photographer.user else 'A Photographer'
        
        notification = Notification(
            user_id=booking.customer_id,
            booking_id=booking.id,
            title='Job Accepted! 🎉',
            message=f'{photographer_name} has accepted your {booking.service_type} booking for {booking.event_date}. Please pay the token amount to confirm your booking.',
            notification_type='job_accepted',
            is_read=False
        )
        db.session.add(notification)
        
        # Send email with PDF invoice to customer
        try:
            customer = User.query.get(booking.customer_id)
            if customer and photographer:
                pdf_bytes = generate_booking_invoice_pdf(booking, customer, photographer)
                send_booking_accepted_email(booking, customer, photographer, pdf_bytes)
        except Exception as e:
            print(f"[EMAIL ERROR] Failed to send booking acceptance email: {e}")
    
    elif data.get('status') == 'declined' and old_status != 'declined':
        photographer = Photographer.query.get(booking.photographer_id)
        photographer_name = photographer.user.first_name if photographer and photographer.user else 'A Photographer'
        
        # --- NEW CODE: Create Refund Request if already paid ---
        if booking.payment_status in ('token_paid', 'paid'):
            refund_amount = booking.token_amount if booking.payment_status == 'token_paid' else booking.total_price
            refund_req = RefundRequest(
                booking_id=booking.id,
                customer_id=booking.customer_id,
                refund_amount=refund_amount,
                reason=f"Photographer declined/cancelled the booking. (By: {photographer_name})",
                status='pending',
                cancelled_by='photographer'
            )
            db.session.add(refund_req)
        # --- END NEW CODE ---
        
        notification = Notification(
            user_id=booking.customer_id,
            booking_id=booking.id,
            title='Job Declined',
            message=f'{photographer_name} has declined your {booking.service_type} booking for {booking.event_date}',
            notification_type='job_declined',
            is_read=False
        )
        db.session.add(notification)
    
    db.session.commit()
    
    return jsonify({
        'message': 'Booking updated',
        'booking': booking.to_dict()
    }), 200


@app.route('/api/bookings/<int:booking_id>/work-complete', methods=['PUT'])
def mark_work_complete(booking_id):
    """Photographer marks work as complete — balance paid, session done"""
    data = request.json or {}
    photographer_id = data.get('photographer_id')

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    if booking.photographer_id != photographer_id:
        return jsonify({'error': 'Unauthorized'}), 403

    if booking.status not in ('token_paid', 'accepted'):
        return jsonify({'error': f'Cannot complete booking with status: {booking.status}'}), 400

    # Record balance payment
    balance_amount = booking.remaining_amount or (booking.total_price - (booking.token_amount or 0))
    if balance_amount > 0:
        balance_payment = Payment(
            booking_id=booking.id,
            customer_id=booking.customer_id,
            photographer_id=booking.photographer_id,
            amount=balance_amount,
            payment_type='remaining',
            payment_method='cash/direct',
            transaction_id=f'BAL-{uuid.uuid4().hex[:10].upper()}',
            status='completed'
        )
        db.session.add(balance_payment)

    # Update booking status
    booking.status = 'completed'
    booking.payment_status = 'paid'

    # Notify customer
    photographer = Photographer.query.get(booking.photographer_id)
    photographer_name = photographer.user.first_name if photographer and photographer.user else 'Your photographer'

    notification = Notification(
        user_id=booking.customer_id,
        booking_id=booking.id,
        title='Session Completed! 🎉📸',
        message=f'{photographer_name} has completed your {booking.service_type} session. Balance payment received. Thank you for choosing BookMyShoot!',
        notification_type='session_completed',
        is_read=False
    )
    db.session.add(notification)

    # Log activity
    activity = ActivityLog(
        user_id=photographer.user_id if photographer else None,
        action='work_completed',
        details=f'Session #{booking.id} ({booking.service_type}) marked as completed. Balance ₹{balance_amount} collected.'
    )
    db.session.add(activity)

    db.session.commit()

    return jsonify({
        'message': 'Work completed successfully',
        'booking': booking.to_dict(),
        'photographer_name': f'{photographer.user.first_name} {photographer.user.last_name}' if photographer and photographer.user else 'Photographer'
    }), 200


# ==================== PAYMENT ROUTES ====================

@app.route('/api/payments', methods=['POST'])
def create_payment():
    """Process a payment"""
    data = request.json
    
    # Validate required fields
    required_fields = ['booking_id', 'customer_id', 'amount', 'payment_type']
    for field in required_fields:
        if not data.get(field):
            return jsonify({'error': f'{field} is required'}), 400
    
    # Verify booking exists
    booking = Booking.query.get(data['booking_id'])
    if not booking:
        return jsonify({'error': 'Booking not found'}), 400
    
    # Generate transaction ID
    transaction_id = f"TXN-{uuid.uuid4().hex[:12].upper()}"
    
    # Create payment record
    payment = Payment(
        booking_id=data['booking_id'],
        customer_id=data['customer_id'],
        photographer_id=booking.photographer_id,
        amount=float(data['amount']),
        payment_type=data['payment_type'],
        payment_method=data.get('payment_method', 'online'),
        transaction_id=transaction_id,
        status='completed'
    )
    
    db.session.add(payment)
    
    # Update booking payment status
    if data['payment_type'] == 'token':
        booking.payment_status = 'token_paid'
    elif data['payment_type'] == 'remaining' or data['payment_type'] == 'full':
        booking.payment_status = 'paid'
    
    db.session.commit()
    
    return jsonify({
        'message': 'Payment successful',
        'payment': payment.to_dict(),
        'booking': booking.to_dict()
    }), 201


@app.route('/api/payments/booking/<int:booking_id>', methods=['GET'])
def get_booking_payments(booking_id):
    """Get all payments for a booking"""
    payments = Payment.query.filter_by(booking_id=booking_id).order_by(Payment.created_at.desc()).all()
    return jsonify({
        'payments': [p.to_dict() for p in payments]
    }), 200


@app.route('/api/bookings/<int:booking_id>/mark_paid', methods=['PUT'])
def mark_booking_paid(booking_id):
    """Mark remaining amount as paid (cash/offline)"""
    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404
        
    booking.payment_status = 'paid'
    booking.status = 'completed' # Auto-complete if fully paid? Or keep separate?
    # User asked: "photographer should update the pending amount that paid"
    
    # Let's log a synthetic payment for record keeping
    transaction_id = f"OFFLINE-{uuid.uuid4().hex[:8].upper()}"
    payment = Payment(
        booking_id=booking.id,
        customer_id=booking.customer_id,
        photographer_id=booking.photographer_id,
        amount=booking.remaining_amount,
        payment_type='remaining',
        payment_method='cash/offline',
        transaction_id=transaction_id,
        status='completed'
    )
    db.session.add(payment)
    db.session.commit()
    
    return jsonify({'message': 'Marked as paid', 'booking': booking.to_dict()}), 200


@app.route('/api/payment/pay-token-offline', methods=['POST'])
def pay_token_offline():
    """Process token payment without Razorpay (demo/offline mode)"""
    data = request.json
    booking_id = data.get('booking_id')
    customer_id = data.get('customer_id')

    if not booking_id or not customer_id:
        return jsonify({'error': 'booking_id and customer_id are required'}), 400

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    if booking.status != 'accepted':
        return jsonify({'error': 'Booking must be accepted by photographer first'}), 400

    if booking.customer_id != int(customer_id):
        return jsonify({'error': 'Unauthorized'}), 403

    if booking.payment_status in ('token_paid', 'paid'):
        return jsonify({'error': 'Token already paid'}), 400

    # Process offline token payment
    transaction_id = f"DEMO-{uuid.uuid4().hex[:10].upper()}"
    payment = Payment(
        booking_id=booking.id,
        customer_id=booking.customer_id,
        photographer_id=booking.photographer_id,
        amount=booking.token_amount,
        payment_type='token',
        payment_method='demo/offline',
        transaction_id=transaction_id,
        status='completed'
    )
    db.session.add(payment)

    booking.payment_status = 'token_paid'
    booking.status = 'token_paid'

    # Notify photographer
    if booking.photographer:
        notification = Notification(
            user_id=booking.photographer.user_id,
            booking_id=booking.id,
            title='Token Payment Received! 💰',
            message=f'Token payment of ₹{booking.token_amount} for {booking.service_type} has been received. Session is confirmed.',
            notification_type='payment_received'
        )
        db.session.add(notification)

    db.session.commit()

    return jsonify({
        'message': 'Token payment successful',
        'payment': payment.to_dict(),
        'booking': booking.to_dict()
    }), 201


# ==================== RAZORPAY PAYMENT ROUTES ====================

@app.route('/api/payment/create-order', methods=['POST'])
def create_razorpay_order():
    """Create a Razorpay order for a booking's token amount"""
    data = request.json
    booking_id = data.get('booking_id')
    customer_id = data.get('customer_id')
    try:
        customer_id = int(customer_id)
    except (TypeError, ValueError):
        return jsonify({'error': 'Invalid customer_id'}), 400

    if not razorpay_client:
        return jsonify({'error': 'Razorpay is not configured. Please set RAZORPAY_KEY_ID and RAZORPAY_KEY_SECRET.'}), 500
    
    booking = Booking.query.get_or_404(booking_id)

    # Ensure booking has been accepted by a photographer
    if booking.status != 'accepted':
        return jsonify({'error': 'Please wait for photographer to accept your booking before making payment'}), 400

    # Ensure only booking owner can initiate token payment
    if booking.customer_id != customer_id:
        return jsonify({'error': 'Unauthorized'}), 403

    if booking.payment_status in ['token_paid', 'paid']:
        return jsonify({'error': 'Token payment already completed for this booking'}), 400
    
    # Amount in paise (multiply by 100)
    amount = int(booking.token_amount * 100)
    
    order_data = {
        'amount': amount,
        'currency': 'INR',
        'receipt': f'receipt_{booking_id}',
        'payment_capture': 1
    }
    
    try:
        razorpay_order = razorpay_client.order.create(data=order_data)
        return jsonify({
            'order_id': razorpay_order['id'],
            'amount': amount,
            'key_id': RAZORPAY_KEY_ID,
            'booking_id': booking_id,
            'customer_name': f"{booking.customer.first_name} {booking.customer.last_name}" if booking.customer else "",
            'customer_email': booking.customer.email if booking.customer else ""
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/payment/verify', methods=['POST'])
def verify_payment():
    """Verify Razorpay payment signature and update status"""
    data = request.json or {}
    razorpay_order_id = data.get('razorpay_order_id')
    razorpay_payment_id = data.get('razorpay_payment_id')
    razorpay_signature = data.get('razorpay_signature')
    booking_id = data.get('booking_id')
    customer_id = data.get('customer_id')
    
    # Validate required Razorpay fields
    if not all([razorpay_order_id, razorpay_payment_id, razorpay_signature, booking_id]):
        return jsonify({'status': 'failure', 'error': 'Missing required payment fields'}), 400
    try:
        customer_id = int(customer_id)
    except (TypeError, ValueError):
        return jsonify({'status': 'failure', 'error': 'Invalid customer_id'}), 400

    if not razorpay_client:
        return jsonify({'status': 'failure', 'error': 'Razorpay is not configured'}), 500
    
    params_dict = {
        'razorpay_order_id': razorpay_order_id,
        'razorpay_payment_id': razorpay_payment_id,
        'razorpay_signature': razorpay_signature
    }
    
    try:
        # Verify the signature
        razorpay_client.utility.verify_payment_signature(params_dict)
        
        booking = Booking.query.get_or_404(booking_id)
        if booking.customer_id != customer_id:
            return jsonify({'status': 'failure', 'error': 'Unauthorized'}), 403

        booking.payment_status = 'token_paid'
        booking.status = 'token_paid'
        
        # Create payment record
        transaction_id = f"RAZOR-{razorpay_payment_id}"
        payment = Payment(
            booking_id=booking.id,
            customer_id=booking.customer_id,
            photographer_id=booking.photographer_id,
            amount=booking.token_amount,
            payment_type='token',
            payment_method='razorpay',
            transaction_id=transaction_id,
            status='completed'
        )
        db.session.add(payment)
        
        # Notify photographer
        if booking.photographer:
            notification = Notification(
                user_id=booking.photographer.user_id,
                booking_id=booking.id,
                title='Token Payment Received! 💰',
                message=f'Token payment for {booking.service_type} has been received. Session is now confirmed.',
                notification_type='payment_received'
            )
            db.session.add(notification)
        
        db.session.commit()
        log_activity(booking.customer_id, 'payment_verified', f"Verified Razorpay payment for booking ID: {booking_id}")
        
        return jsonify({'status': 'success', 'message': 'Payment verified and booking updated'}), 200
    except Exception as e:
        return jsonify({'status': 'failure', 'error': str(e)}), 400


# ==================== FEEDBACK ROUTES ====================

@app.route('/api/feedback', methods=['POST'])
def submit_feedback():
    """Submit feedback for a booking after payment"""
    data = sanitize_data(request.json or {})
    booking_id = data.get('booking_id')
    customer_id = data.get('customer_id')
    rating = data.get('rating')
    comment = data.get('comment', '')

    # Validate feedback fields
    errors = validate_feedback(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    if booking.customer_id != int(customer_id):
        return jsonify({'error': 'Unauthorized'}), 403

    # Check if feedback already exists
    existing = Feedback.query.filter_by(booking_id=booking_id, customer_id=int(customer_id)).first()
    if existing:
        return jsonify({'error': 'Feedback already submitted for this booking'}), 400

    feedback = Feedback(
        booking_id=booking_id,
        customer_id=int(customer_id),
        photographer_id=booking.photographer_id,
        rating=int(rating),
        comment=comment
    )
    db.session.add(feedback)

    # Notify photographer about feedback
    if booking.photographer:
        notification = Notification(
            user_id=booking.photographer.user_id,
            booking_id=booking.id,
            title='New Feedback Received! ⭐',
            message=f'You received a {rating}-star review for your {booking.service_type} session.',
            notification_type='feedback_received'
        )
        db.session.add(notification)

    db.session.commit()

    return jsonify({
        'message': 'Feedback submitted successfully',
        'feedback': feedback.to_dict()
    }), 201


@app.route('/api/feedback/booking/<int:booking_id>', methods=['GET'])
def get_booking_feedback(booking_id):
    """Get feedback for a specific booking"""
    feedback = Feedback.query.filter_by(booking_id=booking_id).first()
    if feedback:
        return jsonify({'feedback': feedback.to_dict()}), 200
    return jsonify({'feedback': None}), 200


# ==================== PHOTOGRAPHER ROUTES ====================

@app.route('/api/photographers', methods=['GET'])
def get_photographers():
    """Get all available photographers"""
    photographers = Photographer.query.filter_by(is_available=True).all()
    return jsonify({
        'photographers': [p.to_dict() for p in photographers]
    }), 200


@app.route('/api/photographers/search', methods=['GET'])
def search_photographers():
    """Search photographers by location, specialty, and availability"""
    location = request.args.get('location', '').strip()
    specialty = request.args.get('specialty', '').strip()
    event_date = request.args.get('event_date', '').strip()
    
    # Start with all available photographers
    all_available = Photographer.query.filter_by(is_available=True).all()
    
    # Score and sort photographers (matching filters get priority, but all are shown)
    scored = []
    for p in all_available:
        score = 0
        # Location match gets +2
        if location and p.location and location.lower() in p.location.lower():
            score += 2
        # Specialty match gets +3
        if specialty and p.specialty and specialty.lower() in p.specialty.lower():
            score += 3
        scored.append((score, p))
    
    # Sort by score descending (best matches first)
    scored.sort(key=lambda x: x[0], reverse=True)
    photographers = [p for _, p in scored]
    
    # If event_date is provided, filter out photographers who have bookings on that date
    if event_date:
        try:
            event_date_obj = datetime.strptime(event_date, '%Y-%m-%d').date()
            booked_photographer_ids = db.session.query(Booking.photographer_id).filter(
                Booking.event_date == event_date_obj,
                Booking.status != 'declined',
                Booking.photographer_id.isnot(None)
            ).distinct().all()
            booked_ids = [b[0] for b in booked_photographer_ids]
            photographers = [p for p in photographers if p.id not in booked_ids]
        except ValueError:
            pass
    
    return jsonify({
        'photographers': [p.to_dict() for p in photographers]
    }), 200


@app.route('/api/admin/photographers', methods=['GET'])
def get_all_photographers_admin():
    """Get all photographers for admin with earnings data"""
    photographers = Photographer.query.all()
    result = []
    for p in photographers:
        p_dict = p.to_dict()
        # Calculate real earnings from completed bookings
        completed = Booking.query.filter_by(photographer_id=p.id, status='completed').all()
        p_dict['total_earnings'] = sum(b.total_price for b in completed)
        p_dict['completed_jobs'] = len(completed)
        p_dict['total_bookings'] = Booking.query.filter_by(photographer_id=p.id).count()
        result.append(p_dict)
    return jsonify({
        'photographers': result
    }), 200


@app.route('/api/admin/settings', methods=['GET'])
def get_admin_settings():
    """Get admin settings including UPI configuration for refunds"""
    from config import UPI_ID, UPI_NAME, UPI_BANK
    return jsonify({
        'upi_id': UPI_ID,
        'upi_name': UPI_NAME,
        'upi_bank': UPI_BANK
    }), 200


@app.route('/api/admin/photographers/<int:photographer_id>/toggle', methods=['PUT'])
def toggle_photographer_status(photographer_id):
    """Enable or disable a photographer (admin only)"""
    photographer = Photographer.query.get(photographer_id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404

    photographer.is_available = not photographer.is_available
    status_label = 'enabled' if photographer.is_available else 'disabled'

    # Also toggle the user account active status
    user = User.query.get(photographer.user_id)
    if user:
        user.is_active = photographer.is_available

    db.session.commit()

    log_activity(None, 'admin_toggle_photographer', f'Photographer #{photographer_id} {status_label}')

    return jsonify({
        'message': f'Photographer {status_label} successfully',
        'photographer': photographer.to_dict()
    }), 200


@app.route('/api/admin/all-portfolios', methods=['GET'])
def get_all_portfolios_admin():
    """Get all portfolio items from all photographers"""
    portfolio_items = PortfolioItem.query.all()
    result = []
    for item in portfolio_items:
        item_dict = item.to_dict()
        if item.photographer and item.photographer.user:
            item_dict['photographer_name'] = f"{item.photographer.user.first_name} {item.photographer.user.last_name}"
        else:
            item_dict['photographer_name'] = f"Photographer #{item.photographer_id}"
        result.append(item_dict)
    return jsonify({'portfolio': result}), 200


@app.route('/api/photographers/<int:photographer_id>', methods=['GET'])
def get_photographer(photographer_id):
    """Get photographer by ID"""
    photographer = Photographer.query.get(photographer_id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404
    
    return jsonify({'photographer': photographer.to_dict()}), 200


@app.route('/api/photographers/<int:photographer_id>', methods=['PUT'])
def update_photographer(photographer_id):
    """Update photographer profile fields (specialty, bio, hourly_rate, location, availability)"""
    photographer = Photographer.query.get(photographer_id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404

    data = sanitize_data(request.json or {})
    
    # Validate photographer update fields
    errors = validate_photographer_update(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400
    
    if 'specialty' in data:
        photographer.specialty = data['specialty']
    if 'bio' in data:
        photographer.bio = data['bio']
    if 'hourly_rate' in data:
        try:
            photographer.hourly_rate = float(data['hourly_rate'])
        except (ValueError, TypeError):
            return jsonify({'error': 'Hourly rate must be a valid number'}), 400
    if 'experience' in data:
        photographer.experience = data['experience']
    if 'location' in data:
        photographer.location = data['location']
    if 'is_available' in data:
        photographer.is_available = bool(data['is_available'])

    db.session.commit()
    return jsonify({'message': 'Profile updated', 'photographer': photographer.to_dict()}), 200


# ==================== PACKAGE CRUD ROUTES ====================

@app.route('/api/packages/photographer/<int:photographer_id>', methods=['GET'])
def get_photographer_packages(photographer_id):
    """Get all packages for a photographer"""
    packages = Package.query.filter_by(photographer_id=photographer_id).order_by(Package.created_at.desc()).all()
    return jsonify({'packages': [p.to_dict() for p in packages]}), 200


@app.route('/api/packages', methods=['POST'])
def create_package():
    """Create a new package for a photographer"""
    data = request.json
    required = ['photographer_id', 'name', 'price']
    for f in required:
        if not data.get(f):
            return jsonify({'error': f'{f} is required'}), 400
    try:
        pkg = Package(
            photographer_id=int(data['photographer_id']),
            name=data['name'],
            description=data.get('description', ''),
            price=float(data['price']),
            category_id=data.get('category_id'),
            image_url=data.get('image_url', '')
        )
        db.session.add(pkg)
        db.session.commit()
        return jsonify({'message': 'Package created', 'package': pkg.to_dict()}), 201
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/packages/<int:package_id>', methods=['PUT'])
def update_package(package_id):
    """Update an existing package"""
    pkg = Package.query.get(package_id)
    if not pkg:
        return jsonify({'error': 'Package not found'}), 404
    data = request.json or {}
    if 'name' in data:
        pkg.name = data['name']
    if 'description' in data:
        pkg.description = data['description']
    if 'price' in data:
        pkg.price = float(data['price'])
    if 'category_id' in data:
        pkg.category_id = data['category_id']
    if 'image_url' in data:
        pkg.image_url = data['image_url']
    db.session.commit()
    return jsonify({'message': 'Package updated', 'package': pkg.to_dict()}), 200


@app.route('/api/packages/<int:package_id>', methods=['DELETE'])
def delete_package(package_id):
    """Delete a package"""
    pkg = Package.query.get(package_id)
    if not pkg:
        return jsonify({'error': 'Package not found'}), 404
    db.session.delete(pkg)
    db.session.commit()
    return jsonify({'message': 'Package deleted'}), 200


@app.route('/api/photographers/revenue/<int:photographer_id>', methods=['GET'])
def get_photographer_revenue(photographer_id):
    """Get revenue statistics for a photographer"""
    completed_bookings = Booking.query.filter_by(photographer_id=photographer_id, status='completed').all()
    all_bookings = Booking.query.filter_by(photographer_id=photographer_id).all()

    total_earnings = sum(b.total_price for b in completed_bookings)
    pending_payments = Booking.query.filter_by(
        photographer_id=photographer_id,
        payment_status='token_paid'
    ).all()
    pending_amount = sum(b.remaining_amount for b in pending_payments)

    # Token revenue already collected
    token_collected = sum(b.token_amount for b in all_bookings if b.payment_status in ('token_paid', 'paid'))

    # Payment history with customer and booking details
    payments = Payment.query.filter_by(photographer_id=photographer_id).order_by(Payment.created_at.desc()).all()
    enriched_payments = []
    for p in payments:
        pd = p.to_dict()
        # Add customer info
        if p.customer:
            pd['customer_name'] = f"{p.customer.first_name} {p.customer.last_name}"
            pd['customer_email'] = p.customer.email
            pd['customer_phone'] = getattr(p.customer, 'phone', '')
        else:
            pd['customer_name'] = 'Unknown'
            pd['customer_email'] = ''
            pd['customer_phone'] = ''
        # Add booking info
        if p.booking:
            pd['service_type'] = p.booking.service_type
            pd['event_date'] = p.booking.event_date.isoformat() if p.booking.event_date else ''
            pd['location'] = p.booking.location or ''
            pd['package_name'] = p.booking.package_name or ''
            pd['booking_status'] = p.booking.status
            pd['total_price'] = p.booking.total_price
        else:
            pd['service_type'] = ''
            pd['event_date'] = ''
            pd['location'] = ''
            pd['package_name'] = ''
            pd['booking_status'] = ''
            pd['total_price'] = 0
        enriched_payments.append(pd)

    # Earnings breakdown by service type
    service_breakdown = {}
    for b in completed_bookings:
        stype = b.service_type or 'Other'
        service_breakdown[stype] = service_breakdown.get(stype, 0) + b.total_price

    # Monthly earnings for the last 12 months
    from datetime import datetime, timedelta
    monthly_earnings = {}
    for b in completed_bookings:
        if b.event_date:
            month_key = b.event_date.strftime('%Y-%m')
            monthly_earnings[month_key] = monthly_earnings.get(month_key, 0) + b.total_price

    # Customer-wise earnings summary
    customer_summary = {}
    for b in completed_bookings:
        cid = b.customer_id
        if cid not in customer_summary:
            cname = 'Unknown'
            if b.customer:
                cname = f"{b.customer.first_name} {b.customer.last_name}"
            customer_summary[cid] = {
                'customer_id': cid,
                'customer_name': cname,
                'total_spent': 0,
                'bookings_count': 0,
                'services': []
            }
        customer_summary[cid]['total_spent'] += b.total_price
        customer_summary[cid]['bookings_count'] += 1
        if b.service_type and b.service_type not in customer_summary[cid]['services']:
            customer_summary[cid]['services'].append(b.service_type)

    return jsonify({
        'total_earnings': total_earnings,
        'token_collected': token_collected,
        'pending_amount': pending_amount,
        'completed_jobs': len(completed_bookings),
        'pending_jobs': len(pending_payments),
        'total_bookings': len(all_bookings),
        'payments': enriched_payments,
        'service_breakdown': service_breakdown,
        'monthly_earnings': monthly_earnings,
        'customer_summary': list(customer_summary.values())
    }), 200

@app.route('/api/portfolio/categories', methods=['GET'])
def get_portfolio_categories():
    """Get available portfolio categories"""
    return jsonify({'categories': PortfolioItem.CATEGORIES}), 200

@app.route('/api/portfolio', methods=['POST'])
def upload_portfolio():
    """Upload portfolio image"""
    if 'image' not in request.files:
        return jsonify({'error': 'No image file'}), 400
        
    file = request.files['image']
    photographer_id = request.form.get('photographer_id')
    caption = request.form.get('caption', '')
    category = request.form.get('category', 'General')
    
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
        
    if file:
        filename = f"{uuid.uuid4().hex}_{file.filename}"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        image_url = f"/uploads/portfolio/{filename}"
        item = PortfolioItem(
            photographer_id=photographer_id,
            image_url=image_url,
            caption=caption,
            category=category
        )
        db.session.add(item)
        db.session.commit()
        
        return jsonify({'message': 'Image uploaded', 'item': item.to_dict()}), 201

@app.route('/api/portfolio/<int:item_id>', methods=['DELETE'])
def delete_portfolio_item(item_id):
    """Delete portfolio image"""
    item = PortfolioItem.query.get(item_id)
    if not item:
        return jsonify({'error': 'Item not found'}), 404
        
    # Optional: Delete file from disk
    try:
        filename = item.image_url.split('/')[-1]
        os.remove(os.path.join(app.config['UPLOAD_FOLDER'], filename))
    except:
        pass
        
    db.session.delete(item)
    db.session.commit()
    return jsonify({'message': 'Item deleted'}), 200

@app.route('/api/stats/customer/<int:customer_id>', methods=['GET'])
def get_customer_stats(customer_id):
    """Get booking statistics for a customer"""
    bookings = Booking.query.filter_by(customer_id=customer_id).all()
    
    total_spent = sum(b.total_price for b in bookings if b.payment_status == 'paid')
    pending_bookings = len([b for b in bookings if b.status in ['pending', 'accepted']])
    completed_bookings = len([b for b in bookings if b.status == 'completed'])
    
    return jsonify({
        'total_spent': total_spent,
        'pending_bookings': pending_bookings,
        'completed_bookings': completed_bookings,
        'total_bookings': len(bookings)
    }), 200


# ==================== ADMIN API ROUTES ====================

@app.route('/api/admin/bookings', methods=['GET'])
def get_admin_bookings():
    """Get all bookings for admin management"""
    status = request.args.get('status')
    if status:
        bookings = Booking.query.filter_by(status=status).order_by(Booking.created_at.desc()).all()
    else:
        bookings = Booking.query.order_by(Booking.created_at.desc()).all()
    return jsonify({
        'bookings': [b.to_dict() for b in bookings]
    }), 200

@app.route('/api/admin/revenue', methods=['GET'])
def get_admin_revenue():
    """Get detailed revenue data for admin"""
    from datetime import datetime, timedelta
    now = datetime.utcnow()

    # All payments
    all_payments = Payment.query.order_by(Payment.created_at.desc()).all()
    completed_payments = [p for p in all_payments if p.status == 'completed']
    total_revenue = sum(p.amount for p in completed_payments)

    # Monthly breakdown (last 12 months)
    monthly = []
    for i in range(11, -1, -1):
        target = now - timedelta(days=i * 30)
        month_start = target.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if target.month == 12:
            month_end = month_start.replace(year=target.year + 1, month=1)
        else:
            month_end = month_start.replace(month=target.month + 1)
        month_payments = [p for p in completed_payments if p.created_at and month_start <= p.created_at < month_end]
        monthly.append({
            'month': month_start.strftime('%b %Y'),
            'label': month_start.strftime('%b'),
            'revenue': sum(p.amount for p in month_payments),
            'count': len(month_payments)
        })

    # By service type
    completed_bookings = Booking.query.filter_by(status='completed').all()
    service_revenue = {}
    for b in completed_bookings:
        svc = b.service_type or 'Other'
        if svc not in service_revenue:
            service_revenue[svc] = {'revenue': 0, 'count': 0}
        service_revenue[svc]['revenue'] += b.total_price or 0
        service_revenue[svc]['count'] += 1

    # Pending vs collected
    all_bookings = Booking.query.all()
    token_collected = sum(b.token_amount for b in all_bookings if b.payment_status in ('token_paid', 'paid'))
    pending_amount = sum(b.remaining_amount for b in all_bookings if b.payment_status == 'token_paid')
    total_paid = sum(b.total_price for b in all_bookings if b.payment_status == 'paid')

    # Enriched payments with customer & photographer names
    enriched_payments = []
    for p in all_payments:
        pd = p.to_dict()
        # Customer info
        cust = User.query.get(p.customer_id) if p.customer_id else None
        pd['customer_name'] = (cust.first_name + ' ' + cust.last_name) if cust else 'Unknown'
        pd['customer_email'] = cust.email if cust else ''
        pd['customer_phone'] = cust.phone if cust else ''
        # Photographer info
        photog = Photographer.query.get(p.photographer_id) if p.photographer_id else None
        if photog and photog.user:
            pd['photographer_name'] = photog.user.first_name + ' ' + photog.user.last_name
        else:
            pd['photographer_name'] = 'Unassigned'
        # Booking details
        booking = Booking.query.get(p.booking_id) if p.booking_id else None
        if booking:
            pd['service_type'] = booking.service_type or ''
            pd['package_name'] = booking.package_name or ''
            pd['event_date'] = booking.event_date.isoformat() if booking.event_date else ''
            pd['location'] = booking.location or ''
            pd['total_price'] = booking.total_price or 0
            pd['booking_status'] = booking.status or ''
        enriched_payments.append(pd)

    # All bookings with customer/photographer for revenue table
    all_bookings_enriched = []
    for b in all_bookings:
        bd = b.to_dict()
        cust = User.query.get(b.customer_id) if b.customer_id else None
        bd['customer_name'] = (cust.first_name + ' ' + cust.last_name) if cust else 'Unknown'
        bd['customer_email'] = cust.email if cust else ''
        bd['customer_phone'] = cust.phone if cust else ''
        photog = Photographer.query.get(b.photographer_id) if b.photographer_id else None
        if photog and photog.user:
            bd['photographer_name'] = photog.user.first_name + ' ' + photog.user.last_name
        else:
            bd['photographer_name'] = 'Unassigned'
        all_bookings_enriched.append(bd)

    return jsonify({
        'total_revenue': total_revenue,
        'token_collected': token_collected,
        'pending_amount': pending_amount,
        'total_paid': total_paid,
        'total_payments': len(completed_payments),
        'monthly': monthly,
        'by_service': service_revenue,
        'recent_payments': enriched_payments,
        'all_bookings': all_bookings_enriched
    }), 200

@app.route('/api/admin/packages', methods=['GET'])
def get_admin_packages():
    """Get all packages across all photographers for admin"""
    packages = Package.query.order_by(Package.created_at.desc()).all()
    result = []
    for p in packages:
        d = p.to_dict()
        photographer = Photographer.query.get(p.photographer_id)
        if photographer and photographer.user:
            d['photographer_name'] = photographer.user.first_name + ' ' + photographer.user.last_name
        else:
            d['photographer_name'] = 'Unknown'
        result.append(d)
    return jsonify({'packages': result}), 200

# ─── SERVICE PACKAGES (Website-level packages on service pages) ─────────────

@app.route('/api/service-packages', methods=['GET'])
def get_service_packages():
    """Get all website service packages, optionally filtered by service_type"""
    service_type = request.args.get('service_type', '')
    query = ServicePackage.query.order_by(ServicePackage.service_type, ServicePackage.sort_order)
    if service_type:
        query = query.filter_by(service_type=service_type)
    packages = query.all()
    return jsonify({'packages': [p.to_dict() for p in packages]}), 200

@app.route('/api/service-packages', methods=['POST'])
def create_service_package():
    """Create a new website service package"""
    import json as json_mod
    data = sanitize_data(request.get_json() or {})
    if not data:
        return jsonify({'error': 'No data provided'}), 400
    required = ['service_type', 'name', 'price']
    for field in required:
        if field not in data:
            return jsonify({'error': f'{field} is required'}), 400
    # Validate price
    err = validate_numeric(data.get('price'), 'Price', min_val=0)
    if err:
        return jsonify({'error': err}), 400
    # Validate name length
    err = validate_text_length(data.get('name'), 'Package name', min_len=2, max_len=100, required=True)
    if err:
        return jsonify({'error': err}), 400
    pkg = ServicePackage(
        service_type=data['service_type'],
        name=data['name'],
        price=float(data['price']),
        features=json_mod.dumps(data.get('features', [])),
        is_featured=data.get('is_featured', False),
        sort_order=data.get('sort_order', 0)
    )
    db.session.add(pkg)
    db.session.commit()
    return jsonify({'message': 'Package created', 'package': pkg.to_dict()}), 201

@app.route('/api/service-packages/<int:pkg_id>', methods=['PUT'])
def update_service_package(pkg_id):
    """Update a website service package"""
    import json as json_mod
    pkg = ServicePackage.query.get(pkg_id)
    if not pkg:
        return jsonify({'error': 'Package not found'}), 404
    data = sanitize_data(request.get_json() or {})
    if not data:
        return jsonify({'error': 'No data provided'}), 400
    if 'price' in data:
        err = validate_numeric(data.get('price'), 'Price', min_val=0)
        if err:
            return jsonify({'error': err}), 400
    if 'name' in data:
        err = validate_text_length(data.get('name'), 'Package name', min_len=2, max_len=100, required=True)
        if err:
            return jsonify({'error': err}), 400
        pkg.name = data['name']
    if 'price' in data:
        pkg.price = float(data['price'])
    if 'features' in data:
        pkg.features = json_mod.dumps(data['features'])
    if 'is_featured' in data:
        pkg.is_featured = data['is_featured']
    if 'sort_order' in data:
        pkg.sort_order = data.get('sort_order', pkg.sort_order)
    if 'service_type' in data:
        pkg.service_type = data['service_type']
    db.session.commit()
    return jsonify({'message': 'Package updated', 'package': pkg.to_dict()}), 200

@app.route('/api/service-packages/<int:pkg_id>', methods=['DELETE'])
def delete_service_package(pkg_id):
    """Delete a website service package"""
    pkg = ServicePackage.query.get(pkg_id)
    if not pkg:
        return jsonify({'error': 'Package not found'}), 404
    db.session.delete(pkg)
    db.session.commit()
    return jsonify({'message': 'Package deleted'}), 200

@app.route('/api/admin/stats', methods=['GET'])
def get_admin_stats():
    """Get platform-wide statistics for admin dashboard"""
    from datetime import datetime, timedelta
    
    # 1. Total Revenue (from completed bookings)
    completed_bookings = Booking.query.filter_by(status='completed').all()
    total_revenue = sum(b.total_price for b in completed_bookings)
    
    # 2. User Stats
    total_users = User.query.count()
    total_customers = User.query.filter_by(user_type='customer').count()
    total_photographers = User.query.filter_by(user_type='photographer').count()
    
    # 3. Booking Stats
    total_bookings = Booking.query.count()
    pending_bookings = Booking.query.filter_by(status='pending').count()
    accepted_bookings = Booking.query.filter_by(status='accepted').count()
    completed_booking_count = len(completed_bookings)
    
    # 4. Recent Bookings (Last 5)
    recent_bookings = Booking.query.order_by(Booking.created_at.desc()).limit(5).all()
    
    # 5. Monthly Revenue (last 6 months from payments)
    monthly_revenue = []
    month_labels = []
    now = datetime.utcnow()
    for i in range(5, -1, -1):
        target = now - timedelta(days=i * 30)
        month_start = target.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if target.month == 12:
            month_end = month_start.replace(year=target.year + 1, month=1)
        else:
            month_end = month_start.replace(month=target.month + 1)
        
        payments = Payment.query.filter(
            Payment.created_at >= month_start,
            Payment.created_at < month_end,
            Payment.status == 'completed'
        ).all()
        monthly_revenue.append(sum(p.amount for p in payments))
        month_labels.append(month_start.strftime('%b'))
    
    # 6. Growth calculation (compare last 2 months)
    if len(monthly_revenue) >= 2 and monthly_revenue[-2] > 0:
        growth = round(((monthly_revenue[-1] - monthly_revenue[-2]) / monthly_revenue[-2]) * 100, 1)
    else:
        growth = 0
    
    return jsonify({
        'revenue': {
            'total': total_revenue,
            'growth': growth,
            'monthly': monthly_revenue,
            'month_labels': month_labels
        },
        'users': {
            'total': total_users,
            'customers': total_customers,
            'photographers': total_photographers
        },
        'bookings': {
            'total': total_bookings,
            'pending': pending_bookings,
            'accepted': accepted_bookings,
            'completed': completed_booking_count
        },
        'recent_activity': [b.to_dict() for b in recent_bookings]
    }), 200

@app.route('/api/admin/users', methods=['GET'])
def get_all_users():
    """Get all users for admin management with pagination"""
    user_type = request.args.get('user_type')
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 50, type=int)
    
    query = User.query
    if user_type:
        query = query.filter_by(user_type=user_type)
    
    # Order by created_at desc and paginate
    pagination = query.order_by(User.created_at.desc()).paginate(page=page, per_page=per_page, error_out=False)
    
    return jsonify({
        'users': [u.to_dict() for u in pagination.items],
        'total': pagination.total,
        'page': page,
        'pages': pagination.pages
    }), 200

@app.route('/api/admin/users', methods=['POST'])
def admin_create_user():
    """Create a new user (Customer or Photographer)"""
    data = sanitize_data(request.json or {})
    try:
        # Validate all fields
        errors = validate_admin_create_user(data)
        if errors:
            return jsonify({'error': errors[0], 'errors': errors}), 400
        
        # Check existing
        if User.query.filter_by(email=data['email']).first():
            return jsonify({'error': 'Email already exists'}), 400
            
        user = User(
            first_name=data['first_name'],
            last_name=data['last_name'],
            email=data['email'],
            phone=data.get('phone', ''),
            user_type=data['user_type']
        )
        user.set_password(data['password'])
        db.session.add(user)
        db.session.flush() # Get user.id without committing
        
        # If photographer, add profile
        if data['user_type'] == 'photographer':
            photographer = Photographer(
                user_id=user.id,
                specialty=data.get('specialty', 'General'),
                hourly_rate=data.get('hourly_rate', 0),
                bio=data.get('bio', '')
            )
            db.session.add(photographer)
            
        db.session.commit()
            
        log_activity(None, 'admin_create_user', f"Created {data['user_type']}: {data['email']}")
        
        return jsonify({'message': 'User created successfully', 'user': user.to_dict()}), 201
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>', methods=['DELETE'])
def delete_user(user_id):
    """Delete a user"""
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
        
    try:
        # Delete related bookings first
        from models import Booking
        Booking.query.filter_by(customer_id=user_id).delete()
        
        # If user is a photographer, delete their photographer profile and portfolio
        if user.user_type == 'photographer':
            from models import Photographer, PortfolioItem
            photographer = Photographer.query.filter_by(user_id=user_id).first()
            if photographer:
                # Delete portfolio items
                PortfolioItem.query.filter_by(photographer_id=photographer.id).delete()
                # Delete photographer
                db.session.delete(photographer)
        
        db.session.delete(user)
        db.session.commit()
        
        log_activity(None, 'admin_delete_user', f"Deleted user ID: {user_id}")
        return jsonify({'message': 'User deleted successfully'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>/toggle', methods=['PUT'])
def toggle_user_status(user_id):
    """Enable or disable a user account (admin only)"""
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404

    if user.user_type == 'admin':
        return jsonify({'error': 'Cannot disable admin accounts'}), 403

    user.is_active = not user.is_active
    status_label = 'enabled' if user.is_active else 'disabled'

    # If photographer, also toggle their availability
    if user.user_type == 'photographer':
        photographer = Photographer.query.filter_by(user_id=user.id).first()
        if photographer:
            photographer.is_available = user.is_active

    db.session.commit()

    log_activity(None, 'admin_toggle_user', f'User #{user_id} ({user.email}) {status_label}')

    return jsonify({
        'message': f'User {status_label} successfully',
        'user': user.to_dict()
    }), 200


@app.route('/api/admin/logs', methods=['GET'])
def get_activity_logs():
    """Get system activity logs"""
    logs = ActivityLog.query.order_by(ActivityLog.timestamp.desc()).limit(100).all()
    return jsonify({'logs': [l.to_dict() for l in logs]}), 200

@app.route('/api/admin/logs/clear', methods=['POST'])
def clear_activity_logs():
    """Clear all activity logs"""
    try:
        ActivityLog.query.delete()
        db.session.commit()
        return jsonify({'message': 'Activity logs cleared successfully'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/reports/download', methods=['GET'])
def download_report():
    """Download system report as CSV"""
    import csv
    import io
    from flask import Response
    
    report_type = request.args.get('type', 'bookings')
    
    output = io.StringIO()
    writer = csv.writer(output)
    
    if report_type == 'bookings':
        writer.writerow(['ID', 'Customer', 'Service', 'Date', 'Amount', 'Status'])
        bookings = Booking.query.all()
        for b in bookings:
            writer.writerow([b.id, b.customer.email if b.customer else 'Unknown', b.service_type, b.event_date, b.total_price, b.status])
            
    elif report_type == 'users':
        writer.writerow(['ID', 'Name', 'Email', 'Type', 'Created At'])
        users = User.query.all()
        for u in users:
            writer.writerow([u.id, f"{u.first_name} {u.last_name}", u.email, u.user_type, u.created_at])
            
    elif report_type == 'logs':
         writer.writerow(['Time', 'Action', 'Details', 'User'])
         logs = ActivityLog.query.all()
         for l in logs:
             writer.writerow([l.timestamp, l.action, l.details, l.user.email if l.user else 'System'])

    return Response(
        output.getvalue(),
        mimetype="text/csv",
        headers={"Content-disposition": f"attachment; filename=report_{report_type}.csv"}
    )


# ==================== PORTFOLIO API ROUTES ====================

@app.route('/api/photographers/<int:id>/portfolio', methods=['GET'])
def get_portfolio(id):
    """Get photographer's portfolio"""
    photographer = Photographer.query.get(id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404
    return jsonify({'portfolio': [p.to_dict() for p in photographer.portfolio]}), 200


@app.route('/api/photographers/<int:id>/portfolio', methods=['POST'])
def upload_portfolio_for_photographer(id):
    """Upload portfolio image for a specific photographer"""
    photographer = Photographer.query.get(id)
    if not photographer:
        return jsonify({'error': 'Photographer not found'}), 404
    if 'image' not in request.files:
        return jsonify({'error': 'No image file'}), 400
    file = request.files['image']
    caption = request.form.get('caption', '')
    category = request.form.get('category', 'General')
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    filename = f"{uuid.uuid4().hex}_{file.filename}"
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    image_url = f"/uploads/portfolio/{filename}"
    item = PortfolioItem(photographer_id=id, image_url=image_url, caption=caption, category=category)
    db.session.add(item)
    db.session.commit()
    return jsonify({'message': 'Image uploaded', 'item': item.to_dict()}), 201


# ==================== POWERPOINT GENERATION ====================

@app.route('/generate', methods=['POST'])
def generate_ppt():
    data = request.json
    topic = data.get('content')
    slides_count = int(data.get('slides', 5))

    if not topic:
        return jsonify({"error": "No content provided"}), 400

    from pptx import Presentation
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Generated using BookMyShoot"

    # Content Slides
    for i in range(1, slides_count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{topic} - Slide {i}"
        content = slide.placeholders[1].text_frame
        content.text = f"Key point {i}"

    filename = f"{uuid.uuid4()}.pptx"
    filepath = os.path.join(PPT_FOLDER, filename)
    prs.save(filepath)

    return jsonify({"download_url": f"/download/{filename}"})


@app.route('/download/<filename>')
def download_file(filename):
    filepath = os.path.join(PPT_FOLDER, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return "File not found", 404


# ==================== NOTIFICATION ROUTES ====================

@app.route('/api/notifications/user/<int:user_id>', methods=['GET'])
def get_user_notifications(user_id):
    """Get all notifications for a user"""
    try:
        # Get query parameters
        limit = request.args.get('limit', 10, type=int)
        unread_only = request.args.get('unread_only', False, type=lambda x: x.lower() == 'true')
        
        query = Notification.query.filter_by(user_id=user_id)
        
        if unread_only:
            query = query.filter_by(is_read=False)
        
        notifications = query.order_by(Notification.created_at.desc()).limit(limit).all()
        
        # Count unread notifications
        unread_count = Notification.query.filter_by(user_id=user_id, is_read=False).count()
        
        return jsonify({
            'notifications': [n.to_dict() for n in notifications],
            'unread_count': unread_count
        }), 200
    except Exception as e:
        return jsonify({'error': f'Error fetching notifications: {e}'}), 500


@app.route('/api/notifications/<int:notification_id>/read', methods=['PUT'])
def mark_notification_read(notification_id):
    """Mark a notification as read"""
    try:
        notification = Notification.query.get(notification_id)
        
        if not notification:
            return jsonify({'error': 'Notification not found'}), 404
        
        notification.is_read = True
        notification.read_at = datetime.utcnow()
        db.session.commit()
        
        return jsonify({
            'message': 'Notification marked as read',
            'notification': notification.to_dict()
        }), 200
    except Exception as e:
        return jsonify({'error': f'Error updating notification: {e}'}), 500


@app.route('/api/notifications/user/<int:user_id>/read-all', methods=['PUT'])
def mark_all_notifications_read(user_id):
    """Mark all notifications as read for a user"""
    try:
        notifications = Notification.query.filter_by(user_id=user_id, is_read=False).all()
        
        for notification in notifications:
            notification.is_read = True
            notification.read_at = datetime.utcnow()
        
        db.session.commit()
        
        return jsonify({
            'message': f'Marked {len(notifications)} notifications as read'
        }), 200
    except Exception as e:
        return jsonify({'error': f'Error updating notifications: {e}'}), 500


@app.route('/api/notifications/<int:notification_id>', methods=['DELETE'])
def delete_notification(notification_id):
    """Delete a notification"""
    try:
        notification = Notification.query.get(notification_id)
        
        if not notification:
            return jsonify({'error': 'Notification not found'}), 404
        
        db.session.delete(notification)
        db.session.commit()
        
        return jsonify({'message': 'Notification deleted'}), 200
    except Exception as e:
        return jsonify({'error': f'Error deleting notification: {e}'}), 500


# ==================== REFUND REQUEST ROUTES ====================

@app.route('/api/refunds', methods=['POST'])
def create_refund_request():
    """Customer requests a refund for a booking"""
    data = request.json or {}
    booking_id = data.get('booking_id')
    customer_id = data.get('customer_id')
    reason = data.get('reason', '').strip()

    if not booking_id or not customer_id or not reason:
        return jsonify({'error': 'booking_id, customer_id, and reason are required'}), 400

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    if booking.customer_id != int(customer_id):
        return jsonify({'error': 'Unauthorized'}), 403

    # Check if refund already requested
    existing = RefundRequest.query.filter_by(booking_id=booking_id, customer_id=int(customer_id)).first()
    if existing:
        return jsonify({'error': 'Refund already requested for this booking'}), 400

    # Only allow refund for bookings with payments made
    if booking.payment_status not in ('token_paid', 'paid'):
        return jsonify({'error': 'No payment found for this booking to refund'}), 400

    # Calculate refund amount based on cancellation timing
    from datetime import date
    days_until_event = (booking.event_date - date.today()).days if booking.event_date else 0

    if days_until_event > 7:
        refund_percent = 100
    elif days_until_event > 3:
        refund_percent = 75
    elif days_until_event > 1:
        refund_percent = 50
    else:
        refund_percent = 0

    paid_amount = booking.token_amount if booking.payment_status == 'token_paid' else booking.total_price
    refund_amount = round(paid_amount * refund_percent / 100, 2)

    if refund_amount <= 0:
        return jsonify({'error': 'Refund not available. Cancellation too close to event date.'}), 400

    refund = RefundRequest(
        booking_id=booking_id,
        customer_id=int(customer_id),
        refund_amount=refund_amount,
        reason=reason
    )
    db.session.add(refund)

    # Notify admins
    admins = User.query.filter_by(user_type='admin').all()
    for admin in admins:
        notif = Notification(
            user_id=admin.id,
            booking_id=booking_id,
            title='Refund Request Submitted',
            message=f'Refund of ₹{refund_amount} requested for booking #{booking_id}. Reason: {reason[:100]}',
            notification_type='refund_request'
        )
        db.session.add(notif)

    db.session.commit()

    return jsonify({
        'message': 'Refund request submitted successfully',
        'refund': refund.to_dict(),
        'refund_percent': refund_percent,
        'days_until_event': days_until_event
    }), 201


@app.route('/api/refunds/customer/<int:customer_id>', methods=['GET'])
def get_customer_refunds(customer_id):
    """Get all refund requests for a customer"""
    refunds = RefundRequest.query.filter_by(customer_id=customer_id).order_by(RefundRequest.created_at.desc()).all()
    return jsonify({'refunds': [r.to_dict() for r in refunds]}), 200


@app.route('/api/refunds/booking/<int:booking_id>', methods=['GET'])
def get_booking_refund(booking_id):
    """Get refund request for a specific booking"""
    refund = RefundRequest.query.filter_by(booking_id=booking_id).first()
    return jsonify({'refund': refund.to_dict() if refund else None}), 200


@app.route('/api/admin/refunds', methods=['GET'])
def get_all_refunds():
    """Get all refund requests for admin"""
    status = request.args.get('status')
    query = RefundRequest.query.order_by(RefundRequest.created_at.desc())
    if status:
        query = query.filter_by(status=status)
    refunds = query.all()
    return jsonify({'refunds': [r.to_dict() for r in refunds]}), 200


@app.route('/api/admin/refunds/<int:refund_id>', methods=['PUT'])
def process_refund(refund_id):
    """Admin approves or rejects a refund"""
    data = request.json or {}
    action = data.get('action')  # 'approve' or 'reject'
    admin_notes = data.get('admin_notes', '')

    if action not in ('approve', 'reject'):
        return jsonify({'error': 'action must be approve or reject'}), 400

    refund = RefundRequest.query.get(refund_id)
    if not refund:
        return jsonify({'error': 'Refund request not found'}), 404

    if refund.status != 'pending':
        return jsonify({'error': 'Refund already processed'}), 400

    # Get UPI details from config
    from config import UPI_ID, UPI_NAME
    
    refund.admin_notes = admin_notes
    refund.resolved_at = datetime.utcnow()

    if action == 'approve':
        refund.status = 'approved'
        # Generate receipt number
        refund.receipt_number = f"REF-{refund.id:06d}-{datetime.utcnow().strftime('%Y%m%d%H%M%S')}"
        # Set UPI details
        refund.upi_id = UPI_ID
        refund.upi_name = UPI_NAME
        refund.refund_method = 'upi'
        
        # Update booking status
        booking = Booking.query.get(refund.booking_id)
        if booking:
            booking.status = 'cancelled'
            booking.payment_status = 'refunded'
        
        # Enhanced notification with UPI details and receipt
        msg_title = 'Refund Approved! 💰'
        msg_text = f'''Your refund request for booking #{refund.booking_id} has been approved!

💵 Refund Amount: ₹{refund.refund_amount:,.2f}
📋 Receipt No: {refund.receipt_number}

UPI Refund Details:
• UPI ID: {UPI_ID}
• Name: {UPI_NAME}

📝 Note: The refund has been processed to your original payment method. Please check your UPI app/bank account within 2-3 business days.

Thank you for using BookMyShoot!'''
    else:
        refund.status = 'rejected'
        msg_title = 'Refund Request Declined'
        msg_text = f'Your refund request for booking #{refund.booking_id} was declined.'

    if admin_notes:
        msg_text += f'\n\nNote from admin: {admin_notes}'

    # Notify customer
    notification = Notification(
        user_id=refund.customer_id,
        booking_id=refund.booking_id,
        title=msg_title,
        message=msg_text,
        notification_type='refund_' + action + 'd'
    )
    db.session.add(notification)

    log_activity(None, f'refund_{action}d', f'Refund #{refund_id} for booking #{refund.booking_id} {action}d. Amount: ₹{refund.refund_amount}. Receipt: {refund.receipt_number}')

    db.session.commit()

    return jsonify({
        'message': f'Refund {action}d successfully',
        'refund': refund.to_dict()
    }), 200


# ==================== PHOTOGRAPHER CANCELLATION ====================

@app.route('/api/bookings/<int:booking_id>/photographer-cancel', methods=['PUT'])
def photographer_cancel_booking(booking_id):
    """Photographer cancels a session. Refund is NOT automatic — a Refund button appears."""
    data = request.json or {}
    photographer_id = data.get('photographer_id')
    reason = data.get('reason', 'Cancelled by photographer').strip()

    if not photographer_id:
        return jsonify({'error': 'photographer_id is required'}), 400
    try:
        photographer_id = int(photographer_id)
    except (TypeError, ValueError):
        return jsonify({'error': 'Invalid photographer_id'}), 400

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    if booking.photographer_id != photographer_id:
        return jsonify({'error': 'Unauthorized — not your booking'}), 403

    # Allow cancellation for accepted or token_paid sessions
    if booking.status not in ('accepted', 'token_paid'):
        return jsonify({'error': f'Cannot cancel a booking with status: {booking.status}'}), 400

    # Check for duplicate cancellation
    existing_refund = RefundRequest.query.filter_by(
        booking_id=booking_id,
        cancelled_by='photographer'
    ).first()
    if existing_refund:
        return jsonify({'error': 'A photographer cancellation already exists for this booking.'}), 400

    # Get the token payment amount
    token_payment = Payment.query.filter_by(
        booking_id=booking_id,
        payment_type='token'
    ).order_by(Payment.created_at.desc()).first()

    refund_amount = 0
    gateway_payment_id = None
    if token_payment and token_payment.status == 'completed':
        refund_amount = round(float(token_payment.amount or booking.token_amount or 0), 2)
        txn_id = token_payment.transaction_id or ''
        if txn_id.startswith('RAZOR-'):
            gateway_payment_id = txn_id.replace('RAZOR-', '', 1)
        elif token_payment.payment_method == 'razorpay' and txn_id.startswith('pay_'):
            gateway_payment_id = txn_id
    elif not token_payment:
        refund_amount = round(float(booking.token_amount or 0), 2)

    # Cancel the booking (but do NOT refund yet — manual refund via button)
    booking.status = 'cancelled'
    booking.cancelled_by = 'photographer'
    # Payment status stays as-is until refund is manually processed

    # Create a PENDING refund request (not auto-approved)
    refund = RefundRequest(
        booking_id=booking_id,
        customer_id=booking.customer_id,
        refund_amount=refund_amount,
        reason=reason,
        status='pending',
        cancelled_by='photographer',
        razorpay_payment_id=gateway_payment_id,
        refund_method='razorpay' if gateway_payment_id else 'upi'
    )
    db.session.add(refund)

    # Notify customer
    photographer = Photographer.query.get(booking.photographer_id)
    photographer_name = f"{photographer.user.first_name} {photographer.user.last_name}" if photographer and photographer.user else 'Your photographer'

    notification = Notification(
        user_id=booking.customer_id,
        booking_id=booking.id,
        title='Session Cancelled by Photographer',
        message=f'{photographer_name} has cancelled your {booking.service_type} session for {booking.event_date}.\n\nReason: {reason}\n\nYour refund of ₹{refund_amount:,.2f} is being processed.',
        notification_type='photographer_cancelled'
    )
    db.session.add(notification)

    # Notify admins
    admins = User.query.filter_by(user_type='admin').all()
    for admin in admins:
        admin_notif = Notification(
            user_id=admin.id,
            booking_id=booking.id,
            title='Photographer Cancelled Session',
            message=f'{photographer_name} cancelled booking #{booking_id} ({booking.service_type}). Refund of ₹{refund_amount:,.2f} is pending.',
            notification_type='photographer_cancelled'
        )
        db.session.add(admin_notif)

    # Log activity
    activity = ActivityLog(
        user_id=photographer.user_id if photographer else None,
        action='photographer_cancellation',
        details=f'Cancelled booking #{booking_id}. Refund of ₹{refund_amount} is pending manual processing.'
    )
    db.session.add(activity)

    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        logger.error(f"DB update failed for photographer cancellation of booking {booking_id}: {e}")
        return jsonify({'error': 'Failed to process cancellation. Please try again.'}), 500

    return jsonify({
        'message': 'Session cancelled successfully. Use the Refund button to process the refund.',
        'booking': booking.to_dict(),
        'refund': refund.to_dict(),
        'refund_amount': refund_amount,
        'refund_status': 'pending'
    }), 200


@app.route('/api/bookings/<int:booking_id>/process-refund', methods=['POST'])
def process_manual_refund(booking_id):
    """Manually process Razorpay refund for a cancelled booking."""
    data = request.json or {}
    photographer_id = data.get('photographer_id')

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({'error': 'Booking not found'}), 404

    # Allow photographer who owns the booking or admin
    if photographer_id:
        try:
            photographer_id = int(photographer_id)
        except (TypeError, ValueError):
            return jsonify({'error': 'Invalid photographer_id'}), 400
        if booking.photographer_id != photographer_id:
            return jsonify({'error': 'Unauthorized — not your booking'}), 403

    # Must be a cancelled or declined booking
    if booking.status not in ['cancelled', 'declined']:
        return jsonify({'error': 'Booking is not cancelled or declined'}), 400

    # Find the pending refund request
    refund = RefundRequest.query.filter_by(
        booking_id=booking_id,
        cancelled_by='photographer',
        status='pending'
    ).first()
    if not refund:
        return jsonify({'error': 'No pending refund found for this booking'}), 404

    # Check for duplicate refund transactions
    token_payment = Payment.query.filter_by(
        booking_id=booking_id,
        payment_type='token'
    ).order_by(Payment.created_at.desc()).first()

    if not token_payment:
        return jsonify({'error': 'No token payment found for refund.'}), 400

    if token_payment.status == 'refunded':
        return jsonify({'error': 'Token payment has already been refunded.'}), 400

    existing_refund_tx = RefundTransaction.query.filter_by(payment_id=token_payment.id).first()
    if existing_refund_tx:
        return jsonify({'error': 'A refund has already been processed for this payment.'}), 400

    # Razorpay refund
    gateway_payment_id = refund.razorpay_payment_id
    if not gateway_payment_id:
        # Try to extract from payment record
        txn_id = token_payment.transaction_id or ''
        if txn_id.startswith('RAZOR-'):
            gateway_payment_id = txn_id.replace('RAZOR-', '', 1)
        elif token_payment.payment_method == 'razorpay' and txn_id.startswith('pay_'):
            gateway_payment_id = txn_id

    if not gateway_payment_id:
        return jsonify({'error': 'No Razorpay payment ID found. Manual bank transfer may be needed.'}), 400

    if not razorpay_client:
        return jsonify({'error': 'Razorpay is not configured on the server.'}), 500

    refund_amount = round(float(refund.refund_amount or token_payment.amount or 0), 2)
    if refund_amount <= 0:
        return jsonify({'error': 'Invalid refund amount.'}), 400

    try:
        refund_data = razorpay_client.payment.refund(
            gateway_payment_id,
            {
                'amount': int(round(refund_amount * 100)),
                'notes': {
                    'booking_id': str(booking_id),
                    'reason': refund.reason or 'Photographer cancellation refund',
                    'initiated_by': 'manual_refund'
                }
            }
        )
    except Exception as e:
        logger.error(f"Razorpay refund failed for booking {booking_id}: {e}")
        return jsonify({'error': 'Failed to process Razorpay refund. Please try again.'}), 502

    razorpay_refund_id = refund_data.get('id')
    if not razorpay_refund_id:
        return jsonify({'error': 'Refund initiated but Razorpay did not return a refund_id.'}), 502

    refund_time = datetime.utcnow()
    receipt_number = f"REF-PC-{booking_id:06d}-{refund_time.strftime('%Y%m%d%H%M%S')}"

    # Update refund request status
    refund.status = 'approved'
    refund.razorpay_refund_id = razorpay_refund_id
    refund.razorpay_payment_id = gateway_payment_id
    refund.resolved_at = refund_time
    refund.receipt_number = receipt_number
    refund.refund_method = 'razorpay'

    # Update booking payment status
    booking.payment_status = 'refunded'
    token_payment.status = 'refunded'

    # Create refund transaction log
    refund_tx = RefundTransaction(
        booking_id=booking_id,
        payment_id=token_payment.id,
        amount=refund_amount,
        refund_status='refunded',
        refund_date=refund_time,
        gateway='razorpay',
        gateway_payment_id=gateway_payment_id,
        gateway_refund_id=razorpay_refund_id,
        notes=refund.reason
    )
    db.session.add(refund_tx)

    # Notify customer
    notification = Notification(
        user_id=booking.customer_id,
        booking_id=booking.id,
        title='Refund Processed',
        message=f'Your refund of ₹{refund_amount:,.2f} for booking #{booking_id} has been processed via Razorpay.\n\nReceipt: {receipt_number}\nRefund ID: {razorpay_refund_id}\n\nPlease allow 2-3 business days for the amount to reflect.',
        notification_type='refund_approved'
    )
    db.session.add(notification)

    # Log activity
    activity = ActivityLog(
        user_id=None,
        action='manual_refund_processed',
        details=f'Refund of ₹{refund_amount} processed for booking #{booking_id}. Razorpay refund: {razorpay_refund_id}.'
    )
    db.session.add(activity)

    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        logger.error(f"Refund processed in Razorpay but DB update failed for booking {booking_id}: {e}")
        return jsonify({
            'error': 'Refund was initiated but database update failed. Please contact support.',
            'razorpay_refund_id': razorpay_refund_id
        }), 500

    return jsonify({
        'message': 'Refund processed successfully!',
        'refund': refund.to_dict(),
        'refund_transaction': refund_tx.to_dict(),
        'refund_amount': refund_amount,
        'razorpay_refund_id': razorpay_refund_id,
        'receipt_number': receipt_number,
        'refund_time': refund_time.isoformat()
    }), 200


@app.route('/api/bookings/<int:booking_id>/cancellation-check', methods=['GET'])
def check_cancellation_eligibility(booking_id):
    """Check if a booking can be cancelled by photographer"""
    photographer_id_raw = request.args.get('photographer_id')
    
    try:
        photographer_id = int(photographer_id_raw) if photographer_id_raw else None
    except (ValueError, TypeError):
        return jsonify({
            'can_cancel': False,
            'reason': f'Invalid photographer_id: {photographer_id_raw}',
            'error': 'Invalid photographer_id'
        }), 400

    if not photographer_id:
        return jsonify({
            'can_cancel': False,
            'reason': 'photographer_id is required',
            'error': 'photographer_id is required'
        }), 400

    booking = Booking.query.get(booking_id)
    if not booking:
        return jsonify({
            'can_cancel': False,
            'reason': f'Booking #{booking_id} not found',
            'error': 'Booking not found'
        }), 404

    if booking.photographer_id != photographer_id:
        # Log this for admin investigation
        logger.warning(f"Unauthorized cancellation check: Photographer {photographer_id} tried checking Booking {booking_id} (owned by {booking.photographer_id})")
        return jsonify({
            'can_cancel': False,
            'reason': 'Unauthorized — this booking is not assigned to you.',
            'error': 'Unauthorized'
        }), 403

    can_cancel = True
    reason = ''

    # Only accepted or token_paid bookings can be cancelled
    if booking.status not in ('accepted', 'token_paid'):
        can_cancel = False
        reason = f'Booking status is "{booking.status}". Only accepted or token-paid sessions can be cancelled.'

    # Prevent duplicate cancellation
    if can_cancel:
        existing_refund = RefundRequest.query.filter_by(booking_id=booking_id, cancelled_by='photographer').first()
        if existing_refund:
            can_cancel = False
            reason = 'A photographer cancellation or refund request already exists for this booking.'

    # Calculate paid amount for display
    paid_amount = 0
    token_payment = Payment.query.filter_by(
        booking_id=booking_id,
        payment_type='token'
    ).order_by(Payment.created_at.desc()).first() if can_cancel else None

    if booking.payment_status in ('token_paid', 'paid'):
        paid_amount = booking.token_amount or (token_payment.amount if token_payment else 0)

    return jsonify({
        'can_cancel': can_cancel,
        'reason': reason,
        'paid_amount': paid_amount,
        'payment_status': booking.payment_status
    }), 200


# ==================== CANCELLATION POLICY ====================

@app.route('/api/cancellation-policy', methods=['GET'])
def get_cancellation_policy():
    """Return cancellation/refund policy details"""
    return jsonify({
        'policy': [
            {'days': '7+', 'refund_percent': 100, 'description': 'Full refund if cancelled 7+ days before event'},
            {'days': '3-7', 'refund_percent': 75, 'description': '75% refund if cancelled 3-7 days before event'},
            {'days': '1-3', 'refund_percent': 50, 'description': '50% refund if cancelled 1-3 days before event'},
            {'days': '0', 'refund_percent': 0, 'description': 'No refund on event day or after'}
        ],
        'photographer_cancellation': {
            'min_hours': 48,
            'refund_percent': 100,
            'description': 'Photographer can cancel only before 48 hours of the event. Full refund will be issued.'
        }
    }), 200


# ==================== SUPPORT TICKET ROUTES ====================

@app.route('/api/support-tickets', methods=['POST'])
def create_support_ticket():
    """Customer creates a support ticket"""
    data = sanitize_data(request.json or {})
    customer_id = data.get('customer_id')
    subject = data.get('subject', '').strip()
    message = data.get('message', '').strip()
    priority = data.get('priority', 'medium')

    if not customer_id or not subject or not message:
        return jsonify({'error': 'customer_id, subject, and message are required'}), 400

    if priority not in ('low', 'medium', 'high'):
        priority = 'medium'

    ticket = SupportTicket(
        customer_id=int(customer_id),
        subject=subject,
        message=message,
        priority=priority
    )
    db.session.add(ticket)

    # Notify admins
    admins = User.query.filter_by(user_type='admin').all()
    for admin in admins:
        notif = Notification(
            user_id=admin.id,
            title='New Support Ticket',
            message=f'Support ticket from customer #{customer_id}: {subject[:80]}',
            notification_type='support_ticket'
        )
        db.session.add(notif)

    db.session.commit()

    return jsonify({
        'message': 'Support ticket created successfully',
        'ticket': ticket.to_dict()
    }), 201


@app.route('/api/support-tickets/customer/<int:customer_id>', methods=['GET'])
def get_customer_tickets(customer_id):
    """Get all support tickets for a customer"""
    tickets = SupportTicket.query.filter_by(customer_id=customer_id).order_by(SupportTicket.created_at.desc()).all()
    return jsonify({'tickets': [t.to_dict() for t in tickets]}), 200


@app.route('/api/admin/support-tickets', methods=['GET'])
def get_all_tickets():
    """Get all support tickets for admin"""
    status = request.args.get('status')
    query = SupportTicket.query.order_by(SupportTicket.created_at.desc())
    if status:
        query = query.filter_by(status=status)
    tickets = query.all()
    return jsonify({'tickets': [t.to_dict() for t in tickets]}), 200


@app.route('/api/admin/support-tickets/<int:ticket_id>', methods=['PUT'])
def update_support_ticket(ticket_id):
    """Admin replies to or updates a support ticket"""
    data = request.json or {}
    ticket = SupportTicket.query.get(ticket_id)
    if not ticket:
        return jsonify({'error': 'Ticket not found'}), 404

    if 'status' in data:
        if data['status'] in ('open', 'in_progress', 'resolved', 'closed'):
            ticket.status = data['status']

    if 'admin_reply' in data:
        ticket.admin_reply = data['admin_reply']

        # Notify customer of reply
        notification = Notification(
            user_id=ticket.customer_id,
            title='Support Reply Received',
            message=f'Admin replied to your ticket: {ticket.subject[:60]}',
            notification_type='support_reply'
        )
        db.session.add(notification)

    db.session.commit()

    return jsonify({
        'message': 'Ticket updated successfully',
        'ticket': ticket.to_dict()
    }), 200


# ==================== ADMIN FEEDBACK MANAGEMENT ====================

@app.route('/api/admin/feedbacks', methods=['GET'])
def get_all_feedbacks():
    """Get all feedbacks for admin management"""
    rating = request.args.get('rating', type=int)
    query = Feedback.query.order_by(Feedback.created_at.desc())
    if rating:
        query = query.filter_by(rating=rating)
    feedbacks = query.all()
    result = []
    for f in feedbacks:
        fd = f.to_dict()
        fd['customer'] = f.customer.to_dict() if f.customer else None
        fd['photographer_name'] = None
        if f.photographer and f.photographer.user:
            fd['photographer_name'] = f'{f.photographer.user.first_name} {f.photographer.user.last_name}'
        fd['service_type'] = f.booking.service_type if f.booking else None
        fd['event_date'] = f.booking.event_date.isoformat() if f.booking and f.booking.event_date else None
        result.append(fd)
    return jsonify({'feedbacks': result}), 200


# ==================== ADVANCED REPORT ROUTES ====================

@app.route('/api/admin/reports/bookings', methods=['GET'])
def get_admin_report_bookings():
    """Advanced booking reports with filters"""
    from datetime import datetime as dt, timedelta

    # Get filter params
    date_from = request.args.get('date_from')
    date_to = request.args.get('date_to')
    service_type = request.args.get('service_type')
    status = request.args.get('status')
    payment_status = request.args.get('payment_status')
    customer_id = request.args.get('customer_id', type=int)
    quick_filter = request.args.get('quick_filter')  # today, yesterday, this_week, this_month

    query = Booking.query

    # Quick date filters
    today = date.today() if 'date' in dir() else __import__('datetime').date.today()
    if quick_filter == 'today':
        query = query.filter(Booking.event_date == today)
    elif quick_filter == 'yesterday':
        query = query.filter(Booking.event_date == today - timedelta(days=1))
    elif quick_filter == 'this_week':
        start_of_week = today - timedelta(days=today.weekday())
        query = query.filter(Booking.event_date >= start_of_week, Booking.event_date <= today)
    elif quick_filter == 'this_month':
        start_of_month = today.replace(day=1)
        query = query.filter(Booking.event_date >= start_of_month, Booking.event_date <= today)
    else:
        # Custom date range
        if date_from:
            try:
                df = dt.strptime(date_from, '%Y-%m-%d').date()
                query = query.filter(Booking.event_date >= df)
            except ValueError:
                pass
        if date_to:
            try:
                dto = dt.strptime(date_to, '%Y-%m-%d').date()
                query = query.filter(Booking.event_date <= dto)
            except ValueError:
                pass

    if service_type:
        query = query.filter(Booking.service_type == service_type)
    if status:
        query = query.filter(Booking.status == status)
    if payment_status:
        query = query.filter(Booking.payment_status == payment_status)
    if customer_id:
        query = query.filter(Booking.customer_id == customer_id)

    bookings = query.order_by(Booking.event_date.desc()).all()

    # Enrich results
    result = []
    for b in bookings:
        bd = b.to_dict()
        cust = User.query.get(b.customer_id) if b.customer_id else None
        bd['customer_name'] = f'{cust.first_name} {cust.last_name}' if cust else 'Unknown'
        bd['customer_email'] = cust.email if cust else ''
        photog = Photographer.query.get(b.photographer_id) if b.photographer_id else None
        if photog and photog.user:
            bd['photographer_name'] = f'{photog.user.first_name} {photog.user.last_name}'
        else:
            bd['photographer_name'] = 'Unassigned'
        result.append(bd)

    # Summary stats
    total_revenue = sum(b.total_price for b in bookings if b.status == 'completed')
    total_cancellations = len([b for b in bookings if b.status == 'cancelled'])
    pending_count = len([b for b in bookings if b.status == 'pending'])

    return jsonify({
        'bookings': result,
        'summary': {
            'total_bookings': len(bookings),
            'total_revenue': total_revenue,
            'total_cancellations': total_cancellations,
            'pending_bookings': pending_count
        }
    }), 200


@app.route('/api/admin/reports/export', methods=['GET'])
def export_admin_report():
    """Export report as CSV with professional format"""
    import csv
    import io
    from flask import Response
    from datetime import datetime as dt, timedelta

    # Reuse filter logic
    date_from = request.args.get('date_from')
    date_to = request.args.get('date_to')
    service_type = request.args.get('service_type')
    status = request.args.get('status')
    payment_status = request.args.get('payment_status')
    customer_id = request.args.get('customer_id', type=int)
    quick_filter = request.args.get('quick_filter')

    query = Booking.query
    today = __import__('datetime').date.today()

    if quick_filter == 'today':
        query = query.filter(Booking.event_date == today)
    elif quick_filter == 'yesterday':
        query = query.filter(Booking.event_date == today - timedelta(days=1))
    elif quick_filter == 'this_week':
        start_of_week = today - timedelta(days=today.weekday())
        query = query.filter(Booking.event_date >= start_of_week, Booking.event_date <= today)
    elif quick_filter == 'this_month':
        start_of_month = today.replace(day=1)
        query = query.filter(Booking.event_date >= start_of_month, Booking.event_date <= today)
    else:
        if date_from:
            try:
                df = dt.strptime(date_from, '%Y-%m-%d').date()
                query = query.filter(Booking.event_date >= df)
            except ValueError:
                pass
        if date_to:
            try:
                dto = dt.strptime(date_to, '%Y-%m-%d').date()
                query = query.filter(Booking.event_date <= dto)
            except ValueError:
                pass

    if service_type:
        query = query.filter(Booking.service_type == service_type)
    if status:
        query = query.filter(Booking.status == status)
    if payment_status:
        query = query.filter(Booking.payment_status == payment_status)
    if customer_id:
        query = query.filter(Booking.customer_id == customer_id)

    bookings = query.order_by(Booking.event_date.desc()).all()

    output = io.StringIO()
    writer = csv.writer(output)

    # Header section
    writer.writerow(['BookMyShoot - Booking Report'])
    writer.writerow([f'Generated: {dt.now().strftime("%d %b %Y, %I:%M %p")}'])
    writer.writerow([f'Total Records: {len(bookings)}'])
    writer.writerow([])

    writer.writerow(['Booking ID', 'Customer Name', 'Customer Email', 'Service', 'Package', 'Event Date', 'Event Time', 'Location', 'Amount', 'Payment Status', 'Booking Status', 'Photographer'])

    for b in bookings:
        cust = User.query.get(b.customer_id) if b.customer_id else None
        cust_name = f'{cust.first_name} {cust.last_name}' if cust else 'Unknown'
        cust_email = cust.email if cust else ''
        photog = Photographer.query.get(b.photographer_id) if b.photographer_id else None
        photog_name = f'{photog.user.first_name} {photog.user.last_name}' if photog and photog.user else 'Unassigned'
        writer.writerow([
            f'BMS-{str(b.id).zfill(4)}', cust_name, cust_email,
            b.service_type, b.package_name or '', b.event_date, b.event_time or '',
            b.location or '', b.total_price, b.payment_status, b.status, photog_name
        ])

    return Response(
        output.getvalue(),
        mimetype='text/csv',
        headers={'Content-Disposition': f'attachment; filename=BookMyShoot_Report_{dt.now().strftime("%Y%m%d")}.csv'}
    )


# ==================== ERROR HANDLERS ====================

def _wants_json():
    """Return True if the client prefers JSON (API calls)."""
    return (request.path.startswith('/api/') or
            request.accept_mimetypes.best == 'application/json')

@app.errorhandler(404)
def not_found_error(error):
    if _wants_json():
        return jsonify({'error': 'Resource not found'}), 404
    return render_template('errors/404.html'), 404

@app.errorhandler(500)
def internal_error(error):
    db.session.rollback()
    if _wants_json():
        return jsonify({'error': 'Internal server error'}), 500
    return render_template('errors/500.html'), 500

@app.errorhandler(403)
def forbidden_error(error):
    if _wants_json():
        return jsonify({'error': 'Access forbidden'}), 403
    return render_template('errors/403.html'), 403

@app.errorhandler(400)
def bad_request_error(error):
    if _wants_json():
        return jsonify({'error': 'Bad request'}), 400
    return render_template('errors/400.html'), 400

<<<<<<< HEAD
# ==================== SERVICE PACKAGES ====================
=======
>>>>>>> fba2361 (Initial commit)

# ─── Contact Form ───────────────────────────────────────────────────────
@app.route('/api/contact', methods=['POST'])
def submit_contact():
    """Receive contact form submissions and save to database"""
    data = sanitize_data(request.json or {})
    if not data:
        return jsonify({'error': 'No data provided'}), 400

    # Validate contact fields
    errors = validate_contact(data)
    if errors:
        return jsonify({'error': errors[0], 'errors': errors}), 400

    try:
        msg = ContactMessage(
            first_name=data['first_name'],
            last_name=data.get('last_name', ''),
            email=data['email'],
            subject=data.get('subject', 'No Subject'),
            message=data['message']
        )
        db.session.add(msg)

        # Send notification to all admin users
        admins = User.query.filter_by(user_type='admin').all()
        for admin in admins:
            notif = Notification(
                user_id=admin.id,
                title='New Contact Message',
                message=f"{data['first_name']} ({data['email']}): {data.get('subject', 'No Subject')}",
                notification_type='contact_message'
            )
            db.session.add(notif)

        db.session.commit()

        log_activity(None, 'contact_form', f"Contact from {data['first_name']} ({data['email']}): {data.get('subject', '')}")

        return jsonify({'message': 'Your message has been received. We will get back to you shortly!'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'Failed to send message: {str(e)}'}), 500


@app.route('/api/contact-messages', methods=['GET'])
def get_contact_messages():
    """Get all contact messages (admin only)"""
    try:
        messages = ContactMessage.query.order_by(ContactMessage.created_at.desc()).all()
        unread = ContactMessage.query.filter_by(is_read=False).count()
        return jsonify({
            'messages': [m.to_dict() for m in messages],
            'total': len(messages),
            'unread': unread
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/contact-messages/<int:msg_id>/read', methods=['PUT'])
def mark_contact_read(msg_id):
    """Mark a contact message as read"""
    try:
        msg = ContactMessage.query.get(msg_id)
        if not msg:
            return jsonify({'error': 'Message not found'}), 404
        msg.is_read = True
        db.session.commit()
        return jsonify({'message': 'Marked as read'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/contact-messages/<int:msg_id>', methods=['DELETE'])
def delete_contact_message(msg_id):
    """Delete a contact message"""
    try:
        msg = ContactMessage.query.get(msg_id)
        if not msg:
            return jsonify({'error': 'Message not found'}), 404
        db.session.delete(msg)
        db.session.commit()
        return jsonify({'message': 'Message deleted'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    PORT = 5000

    # Only check/clear port on initial launch, not on Werkzeug reloader child
    is_reloader_child = os.environ.get('WERKZEUG_RUN_MAIN') == 'true'

    if not is_reloader_child:
        def kill_port(port):
            """Kill any process occupying the port (Windows)."""
            import subprocess
            try:
                result = subprocess.run(
                    ['netstat', '-ano'],
                    capture_output=True, text=True, timeout=5
                )
                for line in result.stdout.splitlines():
                    if f':{port}' in line and 'LISTENING' in line:
                        pid = line.strip().split()[-1]
                        if pid.isdigit() and int(pid) != os.getpid():
                            subprocess.run(
                                ['taskkill', '/PID', pid, '/F'],
                                capture_output=True, timeout=5
                            )
                            print(f'  Killed stale process on port {port} (PID {pid})')
            except Exception:
                pass

        # Check if port is already in use and kill stale processes
        sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        try:
            sock.bind(('127.0.0.1', PORT))
            sock.close()
        except OSError:
            print(f'Port {PORT} is occupied — clearing...')
            sock.close()
            kill_port(PORT)
            import time
            time.sleep(2)

    use_reloader = '--no-reload' not in sys.argv
    app.run(
        debug=True,
        port=PORT,
        use_reloader=use_reloader,
        threaded=True
    )
